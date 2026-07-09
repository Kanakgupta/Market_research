"""
Generate a technical presentation:
  "IEEE 802.15.4, Zigbee & Thread — From User Level to Deep Technical Depth"

Audience: Wireless engineer who already knows Bluetooth and wants to master
802.15.4 / Zigbee / Thread well enough to lead technical discussions and
define products & features.

Run:  python build_deck.py
Output: 802154_Zigbee_Thread_Masterclass.pptx  (same folder)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Theme
# ----------------------------------------------------------------------------
NAVY      = RGBColor(0x0F, 0x24, 0x40)   # deep background
NAVY2     = RGBColor(0x13, 0x30, 0x55)
SLATE     = RGBColor(0x1E, 0x3A, 0x5F)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
FOG       = RGBColor(0xE8, 0xEE, 0xF5)
GREY      = RGBColor(0xB8, 0xC4, 0xD4)
ACCENT    = RGBColor(0x2D, 0x9C, 0xDB)   # bright blue (bridge from BT)
BT_BLUE   = RGBColor(0x00, 0x82, 0xFC)
ZIGBEE    = RGBColor(0xEB, 0x1C, 0x24)   # zigbee brand red
ZIGBEE_D  = RGBColor(0xB5, 0x14, 0x1A)
THREAD    = RGBColor(0x8E, 0x44, 0xC4)   # purple
THREAD_D  = RGBColor(0x6C, 0x2E, 0x9C)
GREEN     = RGBColor(0x27, 0xAE, 0x60)
AMBER     = RGBColor(0xF2, 0xA6, 0x2E)
CARD      = RGBColor(0x17, 0x33, 0x59)

SW = Inches(13.333)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------

def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def bg(slide, color=NAVY):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    _set_fill(r, color)
    r.shadow.inherit = False
    return r


def send_to_back(slide, shape):
    spTree = slide.shapes._spTree
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)


def rect(slide, l, t, w, h, color, line=None, line_w=None, shadow=False, radius=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE,
        l, t, w, h)
    _set_fill(shp, color)
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    if radius is not None:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def txt(slide, l, t, w, h, text, size=18, color=WHITE, bold=False, italic=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Segoe UI", spacing=1.0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return tb


def bullets(slide, l, t, w, h, items, size=16, color=FOG, gap=6, font="Segoe UI"):
    """items = list of (level, text, optional_color, optional_bold)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        lvl = it[0]
        text = it[1]
        c = it[2] if len(it) > 2 and it[2] else color
        b = it[3] if len(it) > 3 else False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        bullet_char = "▸ " if lvl == 0 else ("• " if lvl == 1 else "– ")
        run = p.add_run()
        run.text = bullet_char + text
        run.font.size = Pt(size - lvl * 1)
        run.font.color.rgb = c
        run.font.bold = b
        run.font.name = font
    return tb


def footer(slide, idx, tag=""):
    line = rect(slide, Inches(0.0), Inches(7.05), SW, Pt(2), ACCENT)
    txt(slide, Inches(0.4), Inches(7.06), Inches(9), Inches(0.35),
        "802.15.4 · Zigbee · Thread — Technical Masterclass" + (("   |   " + tag) if tag else ""),
        size=9, color=GREY)
    txt(slide, Inches(12.3), Inches(7.06), Inches(0.9), Inches(0.35),
        str(idx), size=10, color=GREY, align=PP_ALIGN.RIGHT)


def chip(slide, l, t, text, color, tcolor=WHITE, w=None):
    w = w or Inches(1.4)
    c = rect(slide, l, t, w, Inches(0.34), color, radius=0.5)
    tf = c.text_frame
    tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = tcolor; r.font.name = "Segoe UI"
    return c


# ----------------------------------------------------------------------------
# Slide templates
# ----------------------------------------------------------------------------
_slide_no = 0

def content_slide(title, subtitle=None, accent=ACCENT, tag=""):
    global _slide_no
    _slide_no += 1
    s = prs.slides.add_slide(BLANK)
    b = bg(s, NAVY); send_to_back(s, b)
    # header band
    rect(s, 0, 0, SW, Inches(1.15), NAVY2)
    rect(s, 0, Inches(1.15), SW, Pt(3), accent)
    rect(s, 0, 0, Inches(0.16), Inches(1.15), accent)
    txt(s, Inches(0.45), Inches(0.16), Inches(12.3), Inches(0.7), title,
        size=28, color=WHITE, bold=True)
    if subtitle:
        txt(s, Inches(0.47), Inches(0.78), Inches(12.3), Inches(0.35), subtitle,
            size=13, color=GREY)
    footer(s, _slide_no, tag)
    return s


def section_slide(no, title, subtitle, color):
    global _slide_no
    _slide_no += 1
    s = prs.slides.add_slide(BLANK)
    b = bg(s, NAVY); send_to_back(s, b)
    rect(s, 0, Inches(2.55), SW, Inches(2.4), color)
    rect(s, 0, Inches(2.5), SW, Pt(4), WHITE)
    txt(s, Inches(0.8), Inches(1.4), Inches(4), Inches(1.0), no, size=90, color=color, bold=True)
    txt(s, Inches(0.85), Inches(2.75), Inches(11.6), Inches(1.1), title,
        size=40, color=WHITE, bold=True)
    txt(s, Inches(0.9), Inches(3.95), Inches(11.6), Inches(0.8), subtitle,
        size=16, color=FOG)
    footer(s, _slide_no)
    return s


# ============================================================================
# SLIDE 1 — TITLE
# ============================================================================
_slide_no += 1
s = prs.slides.add_slide(BLANK)
b = bg(s, NAVY); send_to_back(s, b)
# decorative bands
rect(s, 0, Inches(0), Inches(0.22), SH, BT_BLUE)
rect(s, Inches(0.22), Inches(0), Inches(0.11), SH, ZIGBEE)
rect(s, Inches(0.33), Inches(0), Inches(0.11), SH, THREAD)
txt(s, Inches(0.9), Inches(0.7), Inches(11), Inches(0.5),
    "WIRELESS TECHNOLOGY MASTERCLASS", size=16, color=ACCENT, bold=True)
txt(s, Inches(0.85), Inches(1.5), Inches(11.6), Inches(2.0),
    "IEEE 802.15.4, Zigbee & Thread", size=52, color=WHITE, bold=True)
txt(s, Inches(0.9), Inches(3.25), Inches(11.4), Inches(0.9),
    "From high-level user concepts to deep protocol internals", size=22, color=FOG)
txt(s, Inches(0.9), Inches(4.0), Inches(11.4), Inches(0.6),
    "A structured bridge for the Bluetooth engineer", size=16, color=GREY, italic=True)
# three tech cards
cards = [
    (Inches(0.9),  "IEEE 802.15.4", "The PHY + MAC foundation", BT_BLUE),
    (Inches(5.1),  "Zigbee",        "Application-rich mesh stack", ZIGBEE),
    (Inches(9.3),  "Thread",        "IPv6 low-power mesh", THREAD),
]
for l, t1, t2, c in cards:
    card = rect(s, l, Inches(4.9), Inches(3.6), Inches(1.55), CARD, line=c, line_w=Pt(1.5), radius=0.06)
    rect(s, l, Inches(4.9), Inches(3.6), Pt(4), c)
    txt(s, l+Inches(0.2), Inches(5.15), Inches(3.2), Inches(0.5), t1, size=20, color=WHITE, bold=True)
    txt(s, l+Inches(0.2), Inches(5.7), Inches(3.2), Inches(0.6), t2, size=13, color=GREY)
txt(s, Inches(0.9), Inches(6.75), Inches(11), Inches(0.4),
    "Learning goal: understand every technical discussion · define products & features · read the market", size=12, color=GREY)

# ============================================================================
# SLIDE 2 — How to read this deck / agenda
# ============================================================================
s = content_slide("How to Read This Deck", "Structured from user-level down to bit-level — you can stop at any depth")
bullets(s, Inches(0.6), Inches(1.5), Inches(6.1), Inches(5.2), [
    (0, "Part 1 — The Landscape", ACCENT, True),
    (1, "Where these fit in the IoT wireless world"),
    (1, "Mental model mapped from your Bluetooth knowledge"),
    (0, "Part 2 — IEEE 802.15.4 (the foundation)", BT_BLUE, True),
    (1, "PHY, MAC, channels, frames, topologies"),
    (0, "Part 3 — Zigbee (the application stack)", ZIGBEE, True),
    (1, "NWK, APS, ZCL/ZDO, clusters, security"),
    (0, "Part 4 — Thread (the IP stack)", THREAD, True),
    (1, "6LoWPAN, IPv6, roles, commissioning, security"),
    (0, "Part 5 — Compare · Matter · Products · Future", GREEN, True),
], size=15, gap=7)
# right callout
rect(s, Inches(7.0), Inches(1.5), Inches(5.7), Inches(5.05), CARD, line=ACCENT, line_w=Pt(1), radius=0.04)
txt(s, Inches(7.25), Inches(1.7), Inches(5.2), Inches(0.5), "The one-sentence mental model", size=17, color=ACCENT, bold=True)
bullets(s, Inches(7.25), Inches(2.35), Inches(5.2), Inches(4.0), [
    (0, "802.15.4 = the radio + link layer", WHITE, True),
    (1, "Like BT's PHY + Link Layer. It moves frames; it does not define apps."),
    (0, "Zigbee = full stack ON TOP of 802.15.4", WHITE, True),
    (1, "Adds mesh routing + a rich application language (clusters)."),
    (0, "Thread = IPv6 stack ON TOP of 802.15.4", WHITE, True),
    (1, "Adds 6LoWPAN + IPv6/UDP mesh. App layer is separate (e.g. Matter)."),
    (0, "Matter = common app layer over Thread/Wi-Fi", AMBER, True),
], size=14, gap=8)

# ============================================================================
# SLIDE 3 — IoT wireless landscape
# ============================================================================
s = content_slide("The IoT Wireless Landscape", "Every technology trades range, data rate, power and topology differently")
# Positioning table
rows = [
    ["Technology", "Band", "Raw rate", "Range", "Topology", "Sweet spot"],
    ["Bluetooth LE", "2.4 GHz", "1–2 Mbps", "10–100 m", "Star / Mesh", "Phone-centric, audio, wearables"],
    ["802.15.4 / Zigbee", "2.4 GHz / sub-GHz", "250 kbps", "10–100 m", "Mesh", "Home/building sensor & control mesh"],
    ["Thread", "2.4 GHz", "250 kbps", "10–100 m", "IPv6 Mesh", "IP-native smart home (Matter)"],
    ["Wi-Fi", "2.4/5/6 GHz", "10 Mbps–Gbps", "20–50 m", "Star (AP)", "High-bandwidth, mains-powered"],
    ["LoRaWAN / NB-IoT", "sub-GHz / cellular", "0.3–50 kbps", "km scale", "Star", "Long-range, tiny payloads"],
]
add_table = s.shapes.add_table(len(rows), len(rows[0]),
    Inches(0.5), Inches(1.45), Inches(12.3), Inches(3.6)).table
widths = [2.4, 2.0, 1.9, 1.5, 1.8, 2.7]
for i, wv in enumerate(widths):
    add_table.columns[i].width = Inches(wv)
for r_i, row in enumerate(rows):
    for c_i, val in enumerate(row):
        cell = add_table.cell(r_i, c_i)
        cell.text = val
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.runs[0]
        run.font.name = "Segoe UI"
        run.font.size = Pt(12 if r_i else 12)
        if r_i == 0:
            run.font.bold = True; run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = SLATE
        else:
            run.font.color.rgb = FOG
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD if r_i % 2 else NAVY2
            if "802.15.4" in val:
                run.font.bold = True; run.font.color.rgb = WHITE
txt(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.4),
    "Key insight: 802.15.4 / Zigbee / Thread live in the SAME 2.4 GHz neighbourhood as Bluetooth and Wi-Fi, "
    "but are engineered for many low-power nodes forming self-healing MESH networks — not for streaming.",
    size=15, color=AMBER)

# ============================================================================
# SLIDE 4 — Bridge from Bluetooth (layer mapping)
# ============================================================================
s = content_slide("Bridge From What You Know: Bluetooth ↔ These Stacks", "Same OSI ideas, different names")
map_rows = [
    ["OSI-ish layer", "Bluetooth LE", "Zigbee", "Thread"],
    ["Application", "GATT profiles / services", "Zigbee Cluster Library (ZCL)", "Matter / CoAP / app of choice"],
    ["App support", "ATT / GATT", "APS (App Support Sublayer)", "UDP + DTLS"],
    ["Network / routing", "L2CAP + BT Mesh (opt.)", "Zigbee NWK (mesh routing)", "IPv6 + 6LoWPAN + MLE routing"],
    ["Link / MAC", "Link Layer", "IEEE 802.15.4 MAC", "IEEE 802.15.4 MAC"],
    ["Physical (radio)", "BLE PHY (GFSK 2.4 GHz)", "IEEE 802.15.4 PHY (O-QPSK)", "IEEE 802.15.4 PHY (O-QPSK)"],
]
tb = s.shapes.add_table(len(map_rows), 4, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.1)).table
for i, wv in enumerate([2.6, 3.2, 3.2, 3.3]):
    tb.columns[i].width = Inches(wv)
for r_i, row in enumerate(map_rows):
    for c_i, val in enumerate(row):
        cell = tb.cell(r_i, c_i)
        cell.text = val
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0]; run.font.name = "Segoe UI"; run.font.size = Pt(12.5)
        if r_i == 0:
            run.font.bold = True; run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = SLATE
        else:
            run.font.color.rgb = FOG
            cell.fill.solid(); cell.fill.fore_color.rgb = CARD if r_i % 2 else NAVY2
            if c_i == 0:
                run.font.bold = True; run.font.color.rgb = WHITE
txt(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.0),
    "Takeaway: Zigbee and Thread SHARE the bottom two layers (802.15.4 PHY+MAC). They diverge above the MAC — "
    "Zigbee builds its own network+app world; Thread runs standard IPv6.", size=15, color=AMBER)

# ============================================================================
# SECTION 2 — 802.15.4
# ============================================================================
section_slide("01", "IEEE 802.15.4 — The Foundation", "The PHY and MAC that both Zigbee and Thread stand on", BT_BLUE)

# 802.15.4 overview
s = content_slide("What IEEE 802.15.4 Actually Is", "A standard for Low-Rate Wireless Personal Area Networks (LR-WPAN)", accent=BT_BLUE)
bullets(s, Inches(0.6), Inches(1.5), Inches(6.1), Inches(5.2), [
    (0, "Defines ONLY the two lowest layers", WHITE, True),
    (1, "PHY (physical radio) and MAC (medium access control)"),
    (1, "Deliberately stops there — upper stacks (Zigbee, Thread, WirelessHART, ISA100, Matter-over-Thread) build on top"),
    (0, "Design goals", WHITE, True),
    (1, "Ultra low power (coin-cell, years of life)"),
    (1, "Low cost, low complexity silicon"),
    (1, "Robust in noisy 2.4 GHz environments"),
    (1, "Support many nodes & mesh-friendly links"),
    (0, "First published 2003; revised 2006, 2011, 2015, 2020", GREY),
], size=15, gap=7)
rect(s, Inches(7.0), Inches(1.5), Inches(5.7), Inches(5.05), CARD, line=BT_BLUE, line_w=Pt(1), radius=0.04)
txt(s, Inches(7.25), Inches(1.7), Inches(5.2), Inches(0.5), "Analogy for a BT engineer", size=17, color=BT_BLUE, bold=True)
bullets(s, Inches(7.25), Inches(2.35), Inches(5.2), Inches(4.2), [
    (0, "802.15.4 ≈ the BLE Controller (PHY + Link Layer)."),
    (0, "It gives you: channels, modulation, framing, addressing, CSMA-CA, ACKs, and a MAC state machine."),
    (0, "It does NOT give you: routing, services, security keys management, or apps."),
    (0, "That gap is exactly why Zigbee and Thread exist.", AMBER, True),
], size=15, gap=10)

# PHY
s = content_slide("802.15.4 PHY — The Radio", "Bands, channels and modulation", accent=BT_BLUE)
bullets(s, Inches(0.55), Inches(1.5), Inches(6.2), Inches(5.2), [
    (0, "Three main band options", WHITE, True),
    (1, "2.4 GHz ISM — global, 16 channels, 250 kbps (most common)"),
    (1, "915 MHz (Americas) — 10 channels, 40 kbps"),
    (1, "868 MHz (Europe) — 1 channel, 20 kbps"),
    (0, "2.4 GHz PHY details", WHITE, True),
    (1, "Channels 11–26, spaced 5 MHz apart"),
    (1, "Modulation: O-QPSK with DSSS (spreading) — robust vs interference"),
    (1, "Chip rate 2 Mchip/s → 250 kbps effective"),
    (0, "Transmit power typically 0 dBm (1 mW); range ~10–100 m", GREY),
    (0, "Coexists with Wi-Fi — channels 15, 20, 25, 26 dodge common Wi-Fi", AMBER),
], size=14.5, gap=6)
# channel diagram
rect(s, Inches(7.0), Inches(1.6), Inches(5.7), Inches(2.4), CARD, line=BT_BLUE, line_w=Pt(1), radius=0.04)
txt(s, Inches(7.2), Inches(1.7), Inches(5.3), Inches(0.4), "2.4 GHz channel map (11–26)", size=13, color=BT_BLUE, bold=True)
for i in range(16):
    x = Inches(7.25 + i * 0.34)
    ch = 11 + i
    col = GREEN if ch in (15, 20, 25, 26) else BT_BLUE
    rect(s, x, Inches(2.25), Inches(0.28), Inches(1.0), col, radius=0.2)
    txt(s, x-Inches(0.02), Inches(3.3), Inches(0.34), Inches(0.35), str(ch), size=8, color=GREY, align=PP_ALIGN.CENTER)
txt(s, Inches(7.2), Inches(3.62), Inches(5.3), Inches(0.35), "Green = quieter vs typical Wi-Fi ch 1/6/11", size=10, color=GREY)
rect(s, Inches(7.0), Inches(4.2), Inches(5.7), Inches(2.35), CARD, line=BT_BLUE, line_w=Pt(1), radius=0.04)
txt(s, Inches(7.2), Inches(4.32), Inches(5.3), Inches(0.4), "vs Bluetooth PHY", size=13, color=BT_BLUE, bold=True)
bullets(s, Inches(7.2), Inches(4.75), Inches(5.35), Inches(1.7), [
    (0, "BLE: 40 ch × 2 MHz, GFSK, adaptive frequency hopping."),
    (0, "802.15.4: 16 ch × 5 MHz, O-QPSK/DSSS, NO hopping — picks one channel and stays (spreading gives robustness)."),
], size=13, gap=6)

# MAC
s = content_slide("802.15.4 MAC — Medium Access & Framing", "How nodes share the air and build links", accent=BT_BLUE)
bullets(s, Inches(0.55), Inches(1.45), Inches(6.2), Inches(5.3), [
    (0, "Channel access: CSMA-CA", WHITE, True),
    (1, "Listen-before-talk; random backoff to avoid collisions"),
    (1, "Optional slotted mode with superframe + beacons (GTS for QoS)"),
    (0, "Four MAC frame types", WHITE, True),
    (1, "Beacon — network timing / presence"),
    (1, "Data — carries the payload"),
    (1, "Acknowledgment — reliable delivery"),
    (1, "MAC command — association, disassociation, etc."),
    (0, "Addressing", WHITE, True),
    (1, "64-bit extended (EUI-64, globally unique) OR"),
    (1, "16-bit short address (assigned on join, saves space)"),
    (1, "16-bit PAN ID identifies the network"),
], size=14, gap=5)
rect(s, Inches(7.0), Inches(1.5), Inches(5.7), Inches(2.55), CARD, line=BT_BLUE, line_w=Pt(1), radius=0.04)
txt(s, Inches(7.2), Inches(1.62), Inches(5.3), Inches(0.4), "MAC frame anatomy", size=13, color=BT_BLUE, bold=True)
parts = [("MHR", SLATE, 1.0), ("Frame\nControl", NAVY2, 1.0), ("Seq#", NAVY2, 0.7),
         ("Addr\nfields", NAVY2, 1.1), ("Payload", BT_BLUE, 1.4), ("FCS", SLATE, 0.7)]
x = 7.2
for name, col, wv in parts:
    rect(s, Inches(x), Inches(2.15), Inches(wv), Inches(1.0), col, radius=0.06)
    txt(s, Inches(x), Inches(2.35), Inches(wv), Inches(0.7), name, size=10, color=WHITE, align=PP_ALIGN.CENTER)
    x += wv + 0.05
txt(s, Inches(7.2), Inches(3.35), Inches(5.3), Inches(0.6), "Max PHY payload 127 bytes → tight! Drives header compression in Thread (6LoWPAN).", size=11, color=AMBER)
rect(s, Inches(7.0), Inches(4.25), Inches(5.7), Inches(2.3), CARD, line=BT_BLUE, line_w=Pt(1), radius=0.04)
txt(s, Inches(7.2), Inches(4.37), Inches(5.3), Inches(0.4), "Device types (MAC level)", size=13, color=BT_BLUE, bold=True)
bullets(s, Inches(7.2), Inches(4.8), Inches(5.35), Inches(1.7), [
    (0, "FFD — Full-Function Device: can route, be coordinator."),
    (0, "RFD — Reduced-Function Device: leaf/end node, sleepy, talks only to its parent."),
], size=13, gap=6)

# 802.15.4 deep RF/MAC capabilities
s = content_slide("802.15.4 RF + Link Layer Capability Deep Dive", "The details that drive real product performance", accent=BT_BLUE)
bullets(s, Inches(0.55), Inches(1.5), Inches(6.3), Inches(5.2), [
    (0, "RF capabilities", WHITE, True),
    (1, "Energy Detect (ED) for channel surveys and adaptive channel planning"),
    (1, "CCA modes (energy, carrier, carrier+energy) feed CSMA-CA decisions"),
    (1, "DSSS processing gain improves resilience in noisy ISM environments"),
    (0, "MAC capabilities", WHITE, True),
    (1, "Unslotted/slotted CSMA-CA with tunable BE/backoff/retry parameters"),
    (1, "ACK request + retry limit settings shape reliability vs latency"),
    (1, "Indirect data + frame-pending bit enables sleepy battery children"),
    (1, "Beacon superframe with optional GTS enables deterministic windows"),
    (0, "Security at link layer", WHITE, True),
    (1, "AES-128-CCM* with frame counters and MIC for replay protection"),
], size=13.5, gap=5)
rect(s, Inches(7.1), Inches(1.5), Inches(5.6), Inches(5.05), CARD, line=BT_BLUE, line_w=Pt(1), radius=0.04)
txt(s, Inches(7.3), Inches(1.62), Inches(5.2), Inches(0.4), "Deployment tuning playbook", size=14, color=BT_BLUE, bold=True)
bullets(s, Inches(7.3), Inches(2.1), Inches(5.2), Inches(4.3), [
    (0, "Channel plan first", WHITE, True),
    (1, "Avoid Wi-Fi overlap; prioritize ch 15/20/25/26 when practical"),
    (0, "Budget airtime", WHITE, True),
    (1, "Keep payloads compact; 127-byte PHY frame cap is absolute"),
    (0, "Tune MAC", WHITE, True),
    (1, "Raise retries for reliability; lower retries for latency-sensitive control"),
    (1, "Tune child poll intervals to trade battery life vs downlink responsiveness"),
    (0, "Design for coexistence", AMBER, True),
    (1, "Measure ED/CCA fail rate + retry counters during field trials"),
], size=13, gap=6)

# Topologies
s = content_slide("802.15.4 Topologies", "From simple star to self-healing mesh", accent=BT_BLUE)
# three topology cards
def node(slide, cx, cy, color=BT_BLUE, r=0.22):
    slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx-r), Inches(cy-r), Inches(2*r), Inches(2*r))
    shp = slide.shapes[-1]; _set_fill(shp, color)
    return (cx, cy)

def link(slide, a, b, color=GREY):
    conn = slide.shapes.add_connector(2, Inches(a[0]), Inches(a[1]), Inches(b[0]), Inches(b[1]))
    conn.line.color.rgb = color; conn.line.width = Pt(1.5)
    conn.shadow.inherit = False

titles = [("Star", 0.6, "One coordinator, all nodes talk to it. Simple, single point of failure."),
          ("Cluster tree", 4.9, "Coordinator + routers form a tree; extends range hop-by-hop."),
          ("Mesh", 9.2, "Routers interconnect; multiple paths → self-healing & resilient.")]
for name, ox, desc in titles:
    rect(s, Inches(ox), Inches(1.5), Inches(3.8), Inches(3.6), CARD, line=BT_BLUE, line_w=Pt(1), radius=0.04)
    txt(s, Inches(ox+0.1), Inches(1.6), Inches(3.6), Inches(0.4), name, size=16, color=BT_BLUE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(ox+0.15), Inches(4.15), Inches(3.5), Inches(0.9), desc, size=11.5, color=FOG, align=PP_ALIGN.CENTER)

# star
c = node(s, 2.5, 2.9, AMBER)
for dx, dy in [(-1.2,-0.4),(1.2,-0.4),(-1.2,0.6),(1.2,0.6),(0,0.9)]:
    p = node(s, 2.5+dx, 2.9+dy, BT_BLUE); link(s, c, p)
# tree
root = node(s, 6.8, 2.2, AMBER)
r1 = node(s, 5.7, 3.0, ACCENT); r2 = node(s, 7.9, 3.0, ACCENT)
link(s, root, r1); link(s, root, r2)
for p in [node(s,5.2,3.8,BT_BLUE), node(s,6.2,3.8,BT_BLUE)]:
    link(s, r1, p)
for p in [node(s,7.5,3.8,BT_BLUE), node(s,8.4,3.8,BT_BLUE)]:
    link(s, r2, p)
# mesh
m = [node(s,10.2,2.3,ACCENT), node(s,12.0,2.3,ACCENT), node(s,10.0,3.4,ACCENT),
     node(s,12.2,3.4,ACCENT), node(s,11.1,2.9,AMBER)]
import itertools
for a, bnode in itertools.combinations(m, 2):
    link(s, a, bnode)
txt(s, Inches(0.6), Inches(5.4), Inches(12.3), Inches(1.3),
    "Zigbee and Thread both use MESH. Mains-powered routers relay for battery 'end devices' that sleep. "
    "Add a device → the mesh grows and finds new paths automatically; remove one → traffic reroutes.",
    size=15, color=AMBER)

# ============================================================================
# SECTION 3 — ZIGBEE
# ============================================================================
section_slide("02", "Zigbee — The Application Stack", "Mesh routing + a rich, interoperable device language on top of 802.15.4", ZIGBEE)

# Zigbee stack
s = content_slide("The Zigbee Stack At A Glance", "Everything above the 802.15.4 MAC is 'Zigbee'", accent=ZIGBEE)
layers = [
    ("Application layer — your device logic + ZCL clusters", ZIGBEE, 0.55),
    ("APS — Application Support Sublayer (binding, endpoints, groups)", ZIGBEE_D, 1.3),
    ("ZDO — Zigbee Device Object (discovery, roles, network mgmt)", SLATE, 2.05),
    ("NWK — Network layer (mesh routing, addressing, security)", ACCENT, 2.8),
    ("IEEE 802.15.4 MAC", BT_BLUE, 3.55),
    ("IEEE 802.15.4 PHY (2.4 GHz / sub-GHz radio)", NAVY2, 4.3),
]
for name, col, top in layers:
    rect(s, Inches(0.7), Inches(1.5+ (top-0.55)), Inches(6.6), Inches(0.62), col, radius=0.05)
    txt(s, Inches(0.85), Inches(1.5+(top-0.55)+0.11), Inches(6.35), Inches(0.42), name, size=12.5, color=WHITE, bold=True)
txt(s, Inches(0.7), Inches(6.35), Inches(6.7), Inches(0.5), "Zigbee = NWK + APS + ZDO + ZCL, standardized by the Connectivity Standards Alliance (CSA).", size=12, color=GREY)
rect(s, Inches(7.6), Inches(1.5), Inches(5.1), Inches(5.05), CARD, line=ZIGBEE, line_w=Pt(1), radius=0.04)
txt(s, Inches(7.8), Inches(1.65), Inches(4.7), Inches(0.4), "What each layer buys you", size=15, color=ZIGBEE, bold=True)
bullets(s, Inches(7.8), Inches(2.2), Inches(4.75), Inches(4.2), [
    (0, "NWK", WHITE, True), (1, "Mesh routing (AODV-style), 16-bit addresses, network key security"),
    (0, "APS", WHITE, True), (1, "Endpoints, bindings, groups; reliable end-to-end delivery"),
    (0, "ZDO", WHITE, True), (1, "Device & service discovery, join/leave, role management"),
    (0, "ZCL", WHITE, True), (1, "The shared 'vocabulary' of attributes & commands that makes brands interoperate"),
], size=13, gap=5)

# Zigbee device roles
s = content_slide("Zigbee Device Roles", "Three roles define the mesh", accent=ZIGBEE)
roles = [
    ("Coordinator (ZC)", ZIGBEE, [
        "Exactly ONE per network",
        "Forms the network, picks channel + PAN ID",
        "Holds trust-center security role",
        "Always powered; often the hub/gateway"]),
    ("Router (ZR)", ACCENT, [
        "Mains-powered, always on",
        "Relays packets → extends the mesh",
        "Can host application endpoints too",
        "e.g. smart bulbs, plugs act as routers"]),
    ("End Device (ZED)", BT_BLUE, [
        "Battery-powered, sleepy (RFD)",
        "Talks only to ONE parent router",
        "Cannot relay for others",
        "e.g. door sensor, remote, thermostat"]),
]
x = 0.6
for name, col, items in roles:
    rect(s, Inches(x), Inches(1.5), Inches(3.95), Inches(4.4), CARD, line=col, line_w=Pt(1.5), radius=0.04)
    rect(s, Inches(x), Inches(1.5), Inches(3.95), Pt(5), col)
    txt(s, Inches(x+0.2), Inches(1.7), Inches(3.6), Inches(0.5), name, size=17, color=col, bold=True)
    bullets(s, Inches(x+0.2), Inches(2.35), Inches(3.6), Inches(3.4),
            [(0, it) for it in items], size=13.5, gap=9)
    x += 4.15
txt(s, Inches(0.6), Inches(6.15), Inches(12.3), Inches(0.9),
    "Mental check: a smart bulb is a Router (always powered, extends mesh); a battery PIR sensor is an End Device "
    "(sleeps to save power). More mains-powered devices = stronger mesh.", size=14, color=AMBER)

# Zigbee application model: endpoints/clusters
s = content_slide("Zigbee Application Model: Endpoints, Clusters, Attributes", "The part that defines PRODUCTS and FEATURES", accent=ZIGBEE)
bullets(s, Inches(0.55), Inches(1.5), Inches(6.3), Inches(5.2), [
    (0, "Node → Endpoints (1–240)", WHITE, True),
    (1, "One physical device can expose many endpoints (e.g. a 2-gang switch = 2 endpoints)"),
    (0, "Endpoint → Clusters", WHITE, True),
    (1, "A cluster = a feature module (On/Off, Level, Color Control, Temperature Measurement…)"),
    (1, "Server side hosts attributes; Client side controls them"),
    (0, "Cluster → Attributes + Commands", WHITE, True),
    (1, "Attribute = a state value (e.g. CurrentLevel = 0–254)"),
    (1, "Command = an action (e.g. Move to Level, Toggle)"),
    (0, "Binding & Groups", WHITE, True),
    (1, "Binding links a switch's client cluster to a bulb's server cluster — direct control"),
    (1, "Groups = one command to many devices (whole-room off)"),
], size=13.5, gap=5)
rect(s, Inches(7.1), Inches(1.5), Inches(5.6), Inches(5.05), CARD, line=ZIGBEE, line_w=Pt(1), radius=0.04)
txt(s, Inches(7.3), Inches(1.62), Inches(5.2), Inches(0.4), "Worked example: a smart dimmer bulb", size=14, color=ZIGBEE, bold=True)
bullets(s, Inches(7.3), Inches(2.1), Inches(5.2), Inches(4.3), [
    (0, "Endpoint 1", WHITE, True),
    (1, "On/Off cluster → attribute OnOff, cmd Toggle"),
    (1, "Level Control → attr CurrentLevel, cmd Move to Level"),
    (1, "Color Control → attrs Hue/Saturation/ColorTemp"),
    (1, "Identify, Basic, Groups, Scenes clusters"),
    (0, "Because clusters are standardized…", AMBER, True),
    (1, "A Philips Hue switch can drive an IKEA bulb — same On/Off & Level cluster semantics. THIS is interoperability."),
], size=13, gap=6)

# Zigbee networking / routing
s = content_slide("Zigbee Networking & Routing", "How a packet crosses the mesh", accent=ZIGBEE)
bullets(s, Inches(0.55), Inches(1.5), Inches(6.3), Inches(5.2), [
    (0, "Addressing", WHITE, True),
    (1, "64-bit IEEE address (permanent) + 16-bit NWK short address (assigned at join)"),
    (0, "Mesh routing (AODV-derived)", WHITE, True),
    (1, "Route discovery floods a Route Request; best path replies with Route Reply"),
    (1, "Routers keep a routing table + route to next hop"),
    (1, "Many-to-one routing optimizes 'everyone → gateway' traffic"),
    (0, "Reliability", WHITE, True),
    (1, "MAC ACK per hop + optional APS ACK end-to-end"),
    (0, "Joining", WHITE, True),
    (1, "Permit-join window; device associates to a parent, gets keys & short address"),
], size=14, gap=6)
# packet path diagram
rect(s, Inches(7.1), Inches(1.6), Inches(5.6), Inches(3.4), CARD, line=ZIGBEE, line_w=Pt(1), radius=0.04)
txt(s, Inches(7.3), Inches(1.72), Inches(5.2), Inches(0.4), "Multi-hop delivery", size=14, color=ZIGBEE, bold=True)
sensor = node(s, 7.7, 3.0, BT_BLUE)
r1 = node(s, 9.1, 2.5, ACCENT); r2 = node(s, 10.5, 3.3, ACCENT)
gw = node(s, 11.9, 2.7, AMBER)
link(s, sensor, r1, ZIGBEE); link(s, r1, r2, ZIGBEE); link(s, r2, gw, ZIGBEE)
txt(s, Inches(7.4), Inches(3.25), Inches(1.2), Inches(0.3), "Sensor", size=9, color=GREY)
txt(s, Inches(8.75), Inches(2.05), Inches(1.2), Inches(0.3), "Router", size=9, color=GREY)
txt(s, Inches(10.15), Inches(3.55), Inches(1.2), Inches(0.3), "Router", size=9, color=GREY)
txt(s, Inches(11.55), Inches(2.25), Inches(1.3), Inches(0.3), "Gateway", size=9, color=GREY)
txt(s, Inches(7.3), Inches(4.2), Inches(5.2), Inches(0.7), "Each hop: CSMA-CA + MAC ACK. Battery sensor sleeps; its parent buffers messages until it polls.", size=11.5, color=FOG)
txt(s, Inches(7.1), Inches(5.3), Inches(5.6), Inches(1.2), "Note: classic Zigbee needs a HUB/gateway to bridge to IP/cloud. This is a key contrast with Thread.", size=13, color=AMBER)

# Zigbee security
s = content_slide("Zigbee Security", "AES-128 everywhere, managed by a Trust Center", accent=ZIGBEE)
bullets(s, Inches(0.6), Inches(1.5), Inches(11.8), Inches(5.2), [
    (0, "Cryptography", WHITE, True),
    (1, "AES-128-CCM* provides encryption + integrity (MIC) at NWK and APS layers"),
    (0, "Key hierarchy", WHITE, True),
    (1, "Network key — shared by all nodes; secures mesh-wide traffic"),
    (1, "Link key — pairwise between two devices; secures APS end-to-end and key transport"),
    (1, "Trust Center Link Key — used when joining to receive the network key"),
    (0, "Trust Center", WHITE, True),
    (1, "Usually the Coordinator; authorizes joins and distributes/rotates keys"),
    (0, "Zigbee 3.0 hardened onboarding", WHITE, True),
    (1, "Install codes (per-device pre-shared secret) → unique link key, no well-known default key"),
    (1, "Unified profile: one interoperable Zigbee 3.0 replaces old Home Automation / Light Link profiles"),
    (0, "Historical weak spot: legacy devices used a well-known default Trust Center key — Zigbee 3.0 install codes fix this.", AMBER, True),
], size=14, gap=5)

# Zigbee capability matrix
s = content_slide("Zigbee Capability Matrix", "What Zigbee gives you beyond raw connectivity", accent=ZIGBEE)
bullets(s, Inches(0.55), Inches(1.5), Inches(6.3), Inches(5.2), [
    (0, "Application capabilities", WHITE, True),
    (1, "Standard ZCL clusters for lighting, metering, HVAC, occupancy and security"),
    (1, "Groups + Scenes for low-latency room-wide automation"),
    (0, "Network capabilities", WHITE, True),
    (1, "Self-healing mesh routing with route discovery and repair"),
    (1, "Many-to-one optimization for 'all nodes to gateway' telemetry"),
    (0, "Operational capabilities", WHITE, True),
    (1, "Strong battery model: sleepy end devices + parent buffering"),
    (1, "Mature certification and massive installed base ecosystem"),
    (0, "Integration capabilities", WHITE, True),
    (1, "Bridge into IP/Matter ecosystems via hub/gateway translation"),
], size=13.5, gap=5)
rect(s, Inches(7.1), Inches(1.5), Inches(5.6), Inches(5.05), CARD, line=ZIGBEE, line_w=Pt(1), radius=0.04)
txt(s, Inches(7.3), Inches(1.62), Inches(5.2), Inches(0.4), "Product-fit guidance", size=14, color=ZIGBEE, bold=True)
bullets(s, Inches(7.3), Inches(2.1), Inches(5.2), Inches(4.3), [
    (0, "Best fit", WHITE, True),
    (1, "Lighting/sensor/control networks with high device count"),
    (1, "Cost-sensitive products needing proven module ecosystem"),
    (0, "Watch-outs", WHITE, True),
    (1, "Gateway architecture quality decides UX and cloud latency"),
    (1, "Channel planning and router density decide mesh stability"),
    (0, "Forward strategy", AMBER, True),
    (1, "Ship Zigbee now, expose to Matter through bridges, keep migration path"),
], size=13, gap=6)

# ============================================================================
# SECTION 4 — THREAD
# ============================================================================
section_slide("03", "Thread — The IPv6 Mesh", "Standard Internet Protocol, down to the tiny sensor", THREAD)

# Thread overview
s = content_slide("What Makes Thread Different", "It is IP — every node has an IPv6 address", accent=THREAD)
bullets(s, Inches(0.6), Inches(1.5), Inches(6.1), Inches(5.2), [
    (0, "Same radio & MAC as Zigbee", WHITE, True),
    (1, "IEEE 802.15.4 (2.4 GHz, 250 kbps, mesh)"),
    (0, "But the network layer is native IPv6", WHITE, True),
    (1, "6LoWPAN compresses IPv6 to fit 127-byte frames"),
    (1, "UDP transport; CoAP/DTLS for messaging + security"),
    (1, "No application layer of its own → app is separate (Matter)"),
    (0, "Designed by Thread Group (Nest/Google origin, 2014)", WHITE, True),
    (0, "Key benefits", THREAD, True),
    (1, "No single hub bottleneck — self-healing, no single point of failure"),
    (1, "Border Router bridges Thread ↔ Wi-Fi/Ethernet at the IP level"),
    (1, "End-to-end IP means cloud/app can address a bulb directly"),
], size=14, gap=6)
rect(s, Inches(7.0), Inches(1.5), Inches(5.7), Inches(5.05), CARD, line=THREAD, line_w=Pt(1), radius=0.04)
txt(s, Inches(7.2), Inches(1.62), Inches(5.3), Inches(0.4), "Thread stack", size=15, color=THREAD, bold=True)
tlayers = [
    ("Application (Matter, CoAP…)", THREAD),
    ("UDP + DTLS", THREAD_D),
    ("IPv6 + routing (MLE, RPL-like)", ACCENT),
    ("6LoWPAN adaptation (compression)", SLATE),
    ("IEEE 802.15.4 MAC", BT_BLUE),
    ("IEEE 802.15.4 PHY", NAVY2),
]
yy = 2.15
for name, col in tlayers:
    rect(s, Inches(7.25), Inches(yy), Inches(5.2), Inches(0.6), col, radius=0.06)
    txt(s, Inches(7.4), Inches(yy+0.12), Inches(4.95), Inches(0.4), name, size=12.5, color=WHITE, bold=True)
    yy += 0.7

# Thread roles
s = content_slide("Thread Device Roles", "Roles are dynamic — the mesh self-organizes", accent=THREAD)
troles = [
    ("Border Router", THREAD, ["Bridges Thread ↔ Wi-Fi/Ethernet (IP)", "Can have MANY per network (redundant)", "e.g. HomePod mini, Nest Hub, Eero"]),
    ("Leader", THREAD_D, ["One elected router that manages the mesh", "Assigns router IDs, holds network data", "If it fails, another is auto-elected"]),
    ("Router", ACCENT, ["Forwards packets, always on", "Keeps mesh connectivity & routes", "REED can promote to Router as needed"]),
    ("End Device (MED/SED)", BT_BLUE, ["Minimal/Sleepy End Device (battery)", "Talks to a parent router only", "Sleeps for years on a coin cell"]),
]
x = 0.55
for name, col, items in troles:
    rect(s, Inches(x), Inches(1.5), Inches(3.02), Inches(4.5), CARD, line=col, line_w=Pt(1.4), radius=0.04)
    rect(s, Inches(x), Inches(1.5), Inches(3.02), Pt(5), col)
    txt(s, Inches(x+0.15), Inches(1.68), Inches(2.75), Inches(0.75), name, size=14.5, color=col, bold=True)
    bullets(s, Inches(x+0.15), Inches(2.5), Inches(2.75), Inches(3.4), [(0, it) for it in items], size=12, gap=9)
    x += 3.18
txt(s, Inches(0.55), Inches(6.2), Inches(12.3), Inches(0.9),
    "Self-healing in action: lose the Leader → re-election; lose a Border Router → another takes over; "
    "no manual reconfiguration. Roles (Leader/Router/REED) shift automatically as devices come and go.", size=13.5, color=AMBER)

# Thread commissioning / working
s = content_slide("Thread: Networking & Commissioning", "How a device joins and how packets flow", accent=THREAD)
bullets(s, Inches(0.55), Inches(1.5), Inches(6.3), Inches(5.2), [
    (0, "Addressing", WHITE, True),
    (1, "Every node gets IPv6 addresses (link-local + mesh-local + optional global)"),
    (1, "RLOC (Routing Locator) encodes position in the mesh"),
    (0, "Routing", WHITE, True),
    (1, "Routers exchange MLE (Mesh Link Establishment) to maintain link quality & routes"),
    (1, "Distance-vector routing among routers; end devices via their parent"),
    (0, "Commissioning (secure onboarding)", WHITE, True),
    (1, "New 'Joiner' + 'Commissioner' (phone/app) authenticate via a PSKd/passphrase"),
    (1, "DTLS handshake (J-PAKE) → device receives network credentials"),
    (1, "Matter uses a QR/setup code that carries the onboarding secret"),
    (0, "Multicast & service discovery via IPv6 — no proprietary binding needed", THREAD),
], size=13.5, gap=5)
rect(s, Inches(7.1), Inches(1.55), Inches(5.6), Inches(5.0), CARD, line=THREAD, line_w=Pt(1), radius=0.04)
txt(s, Inches(7.3), Inches(1.67), Inches(5.2), Inches(0.4), "Border Router: the key idea", size=14, color=THREAD, bold=True)
# simple diagram cloud - wifi - border - thread
rect(s, Inches(7.4), Inches(2.2), Inches(1.5), Inches(0.7), NAVY2, line=GREY, line_w=Pt(1), radius=0.2)
txt(s, Inches(7.4), Inches(2.38), Inches(1.5), Inches(0.4), "Cloud / App", size=11, color=FOG, align=PP_ALIGN.CENTER)
rect(s, Inches(9.4), Inches(2.2), Inches(1.6), Inches(0.7), THREAD_D, radius=0.1)
txt(s, Inches(9.4), Inches(2.3), Inches(1.6), Inches(0.5), "Border\nRouter", size=11, color=WHITE, align=PP_ALIGN.CENTER, bold=True)
conn = s.shapes.add_connector(2, Inches(8.9), Inches(2.55), Inches(9.4), Inches(2.55)); conn.line.color.rgb = GREEN; conn.line.width=Pt(2)
txt(s, Inches(8.6), Inches(2.9), Inches(2.5), Inches(0.3), "Wi-Fi / Ethernet (IPv6)", size=9, color=GREY)
tm = [node(s,8.0,4.2,ACCENT), node(s,9.2,4.6,ACCENT), node(s,10.4,4.1,ACCENT), node(s,9.6,3.6,ACCENT), node(s,11.2,4.7,BT_BLUE)]
for a,bn in itertools.combinations(tm,2):
    link(s, a, bn, THREAD)
brnode=(10.2,2.9)
link(s, brnode, (9.6,3.6), THREAD)
txt(s, Inches(7.4), Inches(5.15), Inches(5.1), Inches(1.2), "Because the Border Router speaks IPv6 both sides, your phone app addresses a Thread bulb like any internet host — no protocol translation gateway.", size=12, color=FOG)

# Thread security
s = content_slide("Thread Security", "Secure by default, at multiple layers", accent=THREAD)
bullets(s, Inches(0.6), Inches(1.5), Inches(11.8), Inches(5.2), [
    (0, "Network-wide MAC security", WHITE, True),
    (1, "All 802.15.4 frames encrypted + authenticated with AES-128-CCM using the network key"),
    (0, "Application/transport security", WHITE, True),
    (1, "DTLS secures commissioning and app messaging (CoAP over UDP)"),
    (0, "Onboarding", WHITE, True),
    (1, "J-PAKE password-authenticated key exchange — device proves it knows the passphrase without sending it"),
    (1, "Only authenticated devices ever receive the network credentials"),
    (0, "Operational hardening", WHITE, True),
    (1, "Key rotation supported; no device can join without the commissioning secret"),
    (1, "Mesh-local addressing keeps internal traffic off the global internet by default"),
    (0, "Result: Thread was designed 'secure-first' — every joined node is authenticated and every frame is encrypted.", THREAD, True),
], size=15, gap=6)

# Thread capability matrix
s = content_slide("Thread Capability Matrix", "IP-native strengths and design trade-offs", accent=THREAD)
bullets(s, Inches(0.55), Inches(1.5), Inches(6.3), Inches(5.2), [
    (0, "Networking capabilities", WHITE, True),
    (1, "Native IPv6 end-to-end addressing per node"),
    (1, "6LoWPAN compression + fragmentation for constrained 802.15.4 frames"),
    (1, "Self-healing role-based mesh (Leader/Router/REED/End Device)"),
    (0, "Platform capabilities", WHITE, True),
    (1, "Multiple Border Routers for redundancy and failover"),
    (1, "No single coordinator bottleneck in steady-state operation"),
    (0, "Application capabilities", WHITE, True),
    (1, "Ideal transport for Matter with standard IP tooling and telemetry"),
    (0, "Security capabilities", WHITE, True),
    (1, "Secure commissioning + authenticated membership by design"),
], size=13.5, gap=5)
rect(s, Inches(7.1), Inches(1.5), Inches(5.6), Inches(5.05), CARD, line=THREAD, line_w=Pt(1), radius=0.04)
txt(s, Inches(7.3), Inches(1.62), Inches(5.2), Inches(0.4), "Product-fit guidance", size=14, color=THREAD, bold=True)
bullets(s, Inches(7.3), Inches(2.1), Inches(5.2), Inches(4.3), [
    (0, "Best fit", WHITE, True),
    (1, "Future-facing smart-home products targeting Matter ecosystems"),
    (1, "Products needing IP-native operations and cloud/app observability"),
    (0, "Watch-outs", WHITE, True),
    (1, "Border Router quality and placement strongly affect onboarding UX"),
    (1, "Still needs careful channel planning in crowded 2.4 GHz spaces"),
    (0, "Forward strategy", AMBER, True),
    (1, "Use multiprotocol SoCs: BLE for commissioning + Thread for mesh data plane"),
], size=13, gap=6)

# ============================================================================
# SECTION 5 — COMPARE / MATTER / PRODUCTS / FUTURE
# ============================================================================
section_slide("04", "Compare · Matter · Market · Future", "Putting it all together for product decisions", GREEN)

# Comparison table
s = content_slide("Zigbee vs Thread vs Bluetooth Mesh", "The decision matrix", accent=GREEN)
cmp_rows = [
    ["Dimension", "Zigbee", "Thread", "Bluetooth Mesh"],
    ["Radio / MAC", "802.15.4", "802.15.4", "BLE (2.4 GHz)"],
    ["Network layer", "Zigbee NWK (proprietary)", "Native IPv6", "Managed-flood (no routing tables)"],
    ["IP-addressable", "No (needs gateway)", "Yes (per node)", "No"],
    ["App layer", "ZCL clusters (built-in)", "Separate (Matter/CoAP)", "BT Mesh models"],
    ["Hub needed", "Yes (coordinator/gateway)", "Border Router (IP bridge)", "Provisioner; phone-friendly"],
    ["Interop story", "Zigbee 3.0 certified", "Matter over Thread", "Mesh models / Matter (emerging)"],
    ["Best for", "Mature sensor/control mesh", "Future-proof IP smart home", "Phone-provisioned, lighting"],
]
tb = s.shapes.add_table(len(cmp_rows), 4, Inches(0.4), Inches(1.4), Inches(12.5), Inches(4.9)).table
for i, wv in enumerate([2.5, 3.4, 3.4, 3.2]):
    tb.columns[i].width = Inches(wv)
for r_i, row in enumerate(cmp_rows):
    for c_i, val in enumerate(row):
        cell = tb.cell(r_i, c_i); cell.text = val
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; run = p.runs[0]
        run.font.name = "Segoe UI"; run.font.size = Pt(12)
        if r_i == 0:
            run.font.bold = True; run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = SLATE
        else:
            run.font.color.rgb = FOG
            cell.fill.solid(); cell.fill.fore_color.rgb = CARD if r_i % 2 else NAVY2
            if c_i == 0:
                run.font.bold = True; run.font.color.rgb = WHITE
            if c_i == 2:
                run.font.color.rgb = RGBColor(0xD9,0xC2,0xF0)

# Matter
s = content_slide("Matter — Why It Changes Everything", "The common application layer that rides on Thread (and Wi-Fi)", accent=AMBER)
bullets(s, Inches(0.55), Inches(1.5), Inches(6.3), Inches(5.2), [
    (0, "What Matter is", WHITE, True),
    (1, "An application-layer standard (CSA) for smart-home interoperability"),
    (1, "Runs OVER Thread or Wi-Fi (and uses BLE only for commissioning)"),
    (0, "Why it matters for you", WHITE, True),
    (1, "Separates 'app language' from 'transport' — Thread carries IP; Matter defines device types"),
    (1, "Backed by Apple, Google, Amazon, Samsung → cross-ecosystem"),
    (0, "Matter data model echoes Zigbee's", WHITE, True),
    (1, "Nodes → Endpoints → Clusters → Attributes/Commands (familiar!)"),
    (1, "In fact ZCL heavily influenced Matter's data model"),
    (0, "Zigbee's clusters concept lives on inside Matter — your Zigbee knowledge transfers directly.", AMBER, True),
], size=13.5, gap=6)
rect(s, Inches(7.1), Inches(1.55), Inches(5.6), Inches(5.0), CARD, line=AMBER, line_w=Pt(1), radius=0.04)
txt(s, Inches(7.3), Inches(1.67), Inches(5.2), Inches(0.4), "Where each layer sits", size=14, color=AMBER, bold=True)
mstack = [("Matter (device types, clusters)", AMBER),
          ("Thread  OR  Wi-Fi (transport)", THREAD),
          ("802.15.4 / IPv6   |   IP", ACCENT),
          ("Radio", NAVY2)]
yy=2.2
for name,col in mstack:
    rect(s, Inches(7.3), Inches(yy), Inches(5.15), Inches(0.75), col, radius=0.06)
    txt(s, Inches(7.45), Inches(yy+0.16), Inches(4.9), Inches(0.45), name, size=13, color=WHITE, bold=True)
    yy+=0.9
txt(s, Inches(7.3), Inches(5.95), Inches(5.2), Inches(0.6), "BLE is used only to commission a Matter device onto Thread/Wi-Fi.", size=11.5, color=GREY)

# Products today
s = content_slide("Products In The Market Today", "Where you'll actually meet these stacks", accent=GREEN)
colA = [
    (0, "Zigbee — shipping at huge scale", ZIGBEE, True),
    (1, "Philips Hue lighting (bridge + bulbs)"),
    (1, "IKEA TRÅDFRI / DIRIGERA"),
    (1, "Samsung SmartThings hub & sensors"),
    (1, "Amazon Echo (4th gen) built-in Zigbee hub"),
    (1, "Aqara, Sonoff, Tuya sensors & switches"),
    (1, "Industrial: metering, building automation"),
]
colB = [
    (0, "Thread — rapidly expanding", THREAD, True),
    (1, "Apple HomePod mini / HomePod (Border Router)"),
    (1, "Google Nest Hub (2nd gen), Nest Wifi"),
    (1, "Amazon Eero routers (Border Router)"),
    (1, "Nanoleaf, Eve, Wemo Thread devices"),
    (1, "Matter-over-Thread bulbs, plugs, sensors"),
    (1, "Nordic / Silicon Labs / TI SoCs inside"),
]
rect(s, Inches(0.55), Inches(1.5), Inches(6.0), Inches(4.9), CARD, line=ZIGBEE, line_w=Pt(1), radius=0.03)
bullets(s, Inches(0.75), Inches(1.7), Inches(5.6), Inches(4.6), colA, size=15, gap=9)
rect(s, Inches(6.75), Inches(1.5), Inches(6.0), Inches(4.9), CARD, line=THREAD, line_w=Pt(1), radius=0.03)
bullets(s, Inches(6.95), Inches(1.7), Inches(5.6), Inches(4.6), colB, size=15, gap=9)
txt(s, Inches(0.55), Inches(6.55), Inches(12.3), Inches(0.6),
    "Silicon reality: many modern SoCs are multiprotocol — the SAME chip can run Zigbee, Thread and BLE, chosen in firmware.", size=13, color=AMBER)

# Future
s = content_slide("Where This Is Heading (Future)", "What to expect over the next few years", accent=GREEN)
bullets(s, Inches(0.6), Inches(1.5), Inches(11.9), Inches(5.2), [
    (0, "Thread + Matter become the default for new smart-home products", THREAD, True),
    (1, "IP-native, multi-admin (one device, many ecosystems), no proprietary hub lock-in"),
    (0, "Zigbee keeps a massive installed base and low-cost advantage", ZIGBEE, True),
    (1, "Continues in lighting, retail, industrial; bridges expose Zigbee devices into Matter"),
    (0, "Convergence via multiprotocol silicon", WHITE, True),
    (1, "One SoC runs Thread + Zigbee + BLE; OTA can even switch protocols in the field"),
    (0, "Thread 1.3+ and beyond", WHITE, True),
    (1, "Native Matter support, better Border Router roaming, larger/credentialed networks, improved reliability"),
    (0, "Bluetooth stays complementary", BT_BLUE, True),
    (1, "BLE for commissioning + phone proximity; Thread/Zigbee for the always-on mesh backbone"),
    (0, "Product strategy takeaway: build on multiprotocol SoCs, ship Matter-over-Thread, keep a Zigbee bridge path for legacy.", AMBER, True),
], size=15, gap=6)

# Cheat sheet / glossary
s = content_slide("One-Page Cheat Sheet", "Fast recall for technical discussions", accent=ACCENT)
g1 = [
    (0, "802.15.4", WHITE, True), (1, "PHY+MAC, 2.4 GHz, 250 kbps, 127-byte frames, mesh-capable"),
    (0, "FFD / RFD", WHITE, True), (1, "Full-function (can route) vs reduced-function (leaf/sleepy)"),
    (0, "PAN ID", WHITE, True), (1, "16-bit network identifier"),
    (0, "Zigbee roles", WHITE, True), (1, "Coordinator / Router / End Device"),
    (0, "Cluster", WHITE, True), (1, "Feature module: attributes (state) + commands (actions)"),
    (0, "Binding", WHITE, True), (1, "Direct client→server link (switch→bulb)"),
]
g2 = [
    (0, "Thread roles", WHITE, True), (1, "Leader / Router / REED / End Device / Border Router"),
    (0, "Border Router", WHITE, True), (1, "Bridges Thread ↔ Wi-Fi/Ethernet at IP level"),
    (0, "6LoWPAN", WHITE, True), (1, "Compresses IPv6 into 802.15.4 frames"),
    (0, "MLE", WHITE, True), (1, "Mesh Link Establishment — maintains Thread routes"),
    (0, "Matter", WHITE, True), (1, "App layer over Thread/Wi-Fi; clusters like Zigbee"),
    (0, "Security", WHITE, True), (1, "AES-128-CCM* everywhere; DTLS/J-PAKE onboarding"),
]
rect(s, Inches(0.5), Inches(1.5), Inches(6.1), Inches(5.0), CARD, radius=0.03)
bullets(s, Inches(0.7), Inches(1.65), Inches(5.7), Inches(4.7), g1, size=13.5, gap=5)
rect(s, Inches(6.75), Inches(1.5), Inches(6.05), Inches(5.0), CARD, radius=0.03)
bullets(s, Inches(6.95), Inches(1.65), Inches(5.7), Inches(4.7), g2, size=13.5, gap=5)

# ============================================================================
# REFERENCES
# ============================================================================
s = content_slide("References & Further Reading", "Sources consulted to build this deck", accent=ACCENT)
refs_left = [
    ("IEEE 802.15.4-2020 standard", "https://standards.ieee.org/ieee/802.15.4/7029/"),
    ("Connectivity Standards Alliance (Zigbee)", "https://csa-iot.org/all-solutions/zigbee/"),
    ("Zigbee Specification / Zigbee 3.0", "https://csa-iot.org/developer-resource/specifications-download-request/"),
    ("Thread Group — Thread technology", "https://www.threadgroup.org/What-is-Thread/Overview"),
    ("Thread Specification & Fundamentals", "https://www.threadgroup.org/support#specifications"),
    ("Matter (CSA) overview", "https://csa-iot.org/all-solutions/matter/"),
    ("OpenThread (Google) documentation", "https://openthread.io/guides/thread-primer"),
    ("RFC 4944 — IPv6 over 802.15.4 (6LoWPAN)", "https://www.rfc-editor.org/rfc/rfc4944"),
]
refs_right = [
    ("RFC 6282 — 6LoWPAN header compression", "https://www.rfc-editor.org/rfc/rfc6282"),
    ("RFC 6550 — RPL routing (context)", "https://www.rfc-editor.org/rfc/rfc6550"),
    ("Silicon Labs — Zigbee & Thread fundamentals", "https://www.silabs.com/wireless/zigbee"),
    ("Nordic Semiconductor — Thread & Matter docs", "https://developer.nordicsemi.com/"),
    ("Texas Instruments — 802.15.4 SimpleLink", "https://www.ti.com/wireless-connectivity/thread/overview.html"),
    ("Espressif — Thread/Zigbee/Matter guides", "https://docs.espressif.com/"),
    ("Zigbee Cluster Library (ZCL) spec", "https://csa-iot.org/developer-resource/specifications-download-request/"),
    ("Wikipedia — Zigbee / Thread / 802.15.4", "https://en.wikipedia.org/wiki/Zigbee"),
]
rect(s, Inches(0.5), Inches(1.45), Inches(6.15), Inches(5.4), CARD, radius=0.02)
rect(s, Inches(6.85), Inches(1.45), Inches(6.0), Inches(5.4), CARD, radius=0.02)
def ref_block(slide, refs, lx):
    tb = slide.shapes.add_textbox(Inches(lx), Inches(1.6), Inches(5.75), Inches(5.15))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for name, url in refs:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_after = Pt(8)
        r = p.add_run(); r.text = name; r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = "Segoe UI"
        p2 = tf.add_paragraph(); p2.space_after = Pt(10)
        r2 = p2.add_run(); r2.text = url; r2.font.size = Pt(11)
        r2.font.color.rgb = ACCENT; r2.font.name = "Consolas"
        try:
            r2.hyperlink.address = url
        except Exception:
            pass
ref_block(s, refs_left, 0.7)
ref_block(s, refs_right, 7.05)

# Closing slide
_slide_no += 1
s = prs.slides.add_slide(BLANK)
b = bg(s, NAVY); send_to_back(s, b)
rect(s, 0, Inches(3.0), SW, Pt(4), ACCENT)
txt(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.0), "You're ready.", size=44, color=WHITE, bold=True)
txt(s, Inches(0.95), Inches(3.3), Inches(11.5), Inches(1.6),
    "802.15.4 gives the radio & link · Zigbee adds mesh + a rich device language · "
    "Thread adds IPv6 · Matter unifies the app layer. You can now follow the discussion, "
    "define products, and read the market.", size=18, color=FOG)
txt(s, Inches(0.95), Inches(5.2), Inches(11), Inches(0.5), "Next step: skim the References slide, then pick one SoC vendor SDK and build a hello-world node.", size=14, color=GREY, italic=True)

# Save
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "802154_Zigbee_Thread_Masterclass.pptx")
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
