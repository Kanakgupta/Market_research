"""
Build the 802.15.4 / Zigbee / Thread masterclass deck using the provided
PowerPoint template.

Input template:
  Presentation templae.pptx

Output:
  802154_Zigbee_Thread_Masterclass_TEMPLATE.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

BASE = Path(__file__).resolve().parent
TEMPLATE = BASE / "Presentation templae.pptx"
OUTPUT = BASE / "802154_Zigbee_Thread_Masterclass_TEMPLATE.pptx"


def remove_all_slides(prs: Presentation) -> None:
    """Remove existing sample slides from the template presentation."""
    slide_ids = list(prs.slides._sldIdLst)
    for sld in slide_ids:
        rid = sld.rId
        prs.part.drop_rel(rid)
        prs.slides._sldIdLst.remove(sld)


def set_title(slide, text: str) -> None:
    if slide.shapes.title is not None:
        slide.shapes.title.text = text


def set_subtitle(slide, text: str) -> None:
    for shape in slide.placeholders:
        try:
            if shape.placeholder_format.type == 4:  # SUBTITLE
                shape.text = text
                return
        except Exception:
            continue


def object_placeholders(slide):
    objs = []
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.type == 7:  # OBJECT
                objs.append((ph.placeholder_format.idx, ph))
        except Exception:
            continue
    objs.sort(key=lambda x: x[0])
    return [p for _, p in objs]


def write_bullets(shape, items):
    """
    items: list of tuples -> (level, text)
    """
    tf = shape.text_frame
    tf.clear()
    first = True
    for level, text in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.text = text


def write_lines(shape, lines):
    tf = shape.text_frame
    tf.clear()
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line


def add_title_subtitle(prs, title, subtitle, layout=0):
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    set_title(slide, title)
    set_subtitle(slide, subtitle)
    return slide


def add_title_content(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[4])  # Title | content
    set_title(slide, title)
    content_ph = object_placeholders(slide)[0]
    write_bullets(content_ph, bullets)
    return slide


def add_two_content(prs, title, left_lines, right_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 2 contents
    set_title(slide, title)
    ph = object_placeholders(slide)
    write_lines(ph[0], left_lines)
    write_lines(ph[1], right_lines)
    return slide


def add_three_columns(prs, title, left_lines, center_lines, right_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[11])
    set_title(slide, title)
    ph = object_placeholders(slide)
    write_lines(ph[0], left_lines)
    write_lines(ph[1], center_lines)
    write_lines(ph[2], right_lines)
    return slide


def add_four_contents(prs, title, tl, tr, bl, br):
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    set_title(slide, title)
    ph = object_placeholders(slide)
    write_lines(ph[0], tl)
    write_lines(ph[1], tr)
    write_lines(ph[2], bl)
    write_lines(ph[3], br)
    return slide


def add_section(prs, title, layout=6):
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    set_title(slide, title)
    return slide


def add_final(prs, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[17])
    tb = slide.shapes.add_textbox(prs.slide_width * 0.08, prs.slide_height * 0.25,
                                  prs.slide_width * 0.84, prs.slide_height * 0.5)
    tf = tb.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if i == 0:
            p.alignment = PP_ALIGN.CENTER
    return slide


def build():
    prs = Presentation(str(TEMPLATE))
    remove_all_slides(prs)

    # 1) Cover + learning path
    add_title_subtitle(
        prs,
        "IEEE 802.15.4, Zigbee & Thread",
        "From user-level understanding to deep technical depth | Built for a Bluetooth expert",
        layout=0,
    )

    add_title_content(prs, "What You Will Get From This Deck", [
        (0, "Understand all technical discussions around 802.15.4, Zigbee and Thread"),
        (0, "Build a mental model from high-level concepts down to packet-level behavior"),
        (0, "Define products and features with confidence"),
        (0, "Compare market solutions and choose the right stack for future products"),
        (0, "Connect the dots between Zigbee, Thread, Bluetooth and Matter"),
    ])

    add_title_content(prs, "Agenda", [
        (0, "1. Wireless landscape and Bluetooth bridge"),
        (0, "2. IEEE 802.15.4 foundation: PHY, MAC, channels, topology"),
        (0, "3. Zigbee stack: roles, routing, clusters, security"),
        (0, "4. Thread stack: IPv6 mesh, roles, commissioning, security"),
        (0, "5. Compare technologies, Matter, products in market, future direction"),
        (0, "6. References"),
    ])

    # 2) Landscape
    add_section(prs, "Part 1: The Landscape", layout=6)

    add_two_content(
        prs,
        "Where These Technologies Sit",
        [
            "Bluetooth LE:",
            "- High data rate vs 802.15.4" ,
            "- Star topology by default" ,
            "- Phone-centric use cases" ,
            "",
            "Wi-Fi:",
            "- Very high throughput" ,
            "- High power, AP based" ,
            "",
            "LoRa/NB-IoT:",
            "- Very long range" ,
            "- Low payload and latency trade-offs",
        ],
        [
            "802.15.4 / Zigbee / Thread:",
            "- 2.4 GHz low-power mesh backbone" ,
            "- 250 kbps raw PHY (2.4 GHz)" ,
            "- Many nodes, low duty cycle" ,
            "",
            "Sweet spot:",
            "- Sensors" ,
            "- Switches and lighting" ,
            "- Building/home automation" ,
            "- Multi-year battery operation",
        ],
    )

    add_title_content(prs, "Bluetooth Mental Bridge", [
        (0, "802.15.4 is roughly analogous to BLE PHY + Link Layer"),
        (0, "Zigbee and Thread share the same 802.15.4 radio/MAC foundation"),
        (0, "Zigbee adds its own networking + application framework (clusters)"),
        (0, "Thread adds IPv6 networking (6LoWPAN) and uses app layers like Matter"),
        (0, "Matter is app-layer interoperability over Thread/Wi-Fi (BLE for commissioning)"),
    ])

    # 3) 802.15.4
    add_section(prs, "Part 2: IEEE 802.15.4 Foundation", layout=6)

    add_title_content(prs, "802.15.4: Scope and Design Goals", [
        (0, "Defines only PHY and MAC for LR-WPAN"),
        (0, "Low power, low complexity, low cost"),
        (0, "Reliable short-range links with many nodes"),
        (0, "Supports star/tree/mesh-capable deployments"),
        (0, "Upper layers are intentionally left to stacks like Zigbee and Thread"),
    ])

    add_two_content(
        prs,
        "802.15.4 PHY Details",
        [
            "Bands and rates:",
            "- 2.4 GHz global: 16 channels (11-26), 250 kbps",
            "- 915 MHz: 10 channels, 40 kbps",
            "- 868 MHz: 1 channel, 20 kbps",
            "",
            "2.4 GHz specifics:",
            "- O-QPSK with DSSS",
            "- 5 MHz channel spacing",
            "- Typical practical range 10-100 m",
        ],
        [
            "Coexistence notes:",
            "- Shares ISM spectrum with Wi-Fi and Bluetooth",
            "- Channel planning is critical in dense deployments",
            "- Channels 15/20/25/26 often preferred vs Wi-Fi overlap",
            "",
            "Compared with BLE:",
            "- BLE uses AFH across 2 MHz channels",
            "- 802.15.4 typically stays on a selected channel",
        ],
    )

    add_title_content(prs, "802.15.4 MAC and Frames", [
        (0, "CSMA-CA medium access (listen-before-talk + random backoff)"),
        (0, "Optional beacon/superframe mode with GTS support"),
        (0, "Frame types: Beacon, Data, Acknowledgment, MAC Command"),
        (0, "Addressing: 64-bit EUI-64 and optional 16-bit short addresses"),
        (0, "PAN ID identifies the network"),
        (0, "Max PHY frame size is 127 bytes (major design constraint)"),
    ])

    add_three_columns(
        prs,
        "Topology Model",
        [
            "Star",
            "- One coordinator",
            "- Simple deployment",
            "- Coordinator is bottleneck",
        ],
        [
            "Tree / Cluster Tree",
            "- Hierarchical extension",
            "- Predictable structure",
            "- Less resilient than full mesh",
        ],
        [
            "Mesh",
            "- Multiple paths",
            "- Self-healing",
            "- Preferred for real products",
        ],
    )

    # 4) Zigbee
    add_section(prs, "Part 3: Zigbee", layout=6)

    add_title_content(prs, "Zigbee Stack", [
        (0, "Built on 802.15.4 PHY/MAC"),
        (0, "Adds NWK: mesh routing, addressing, key-based network security"),
        (0, "Adds APS: endpoints, binding, groups, app-level transport"),
        (0, "Adds ZDO: discovery, roles, network management"),
        (0, "Adds ZCL: standardized clusters, attributes, commands"),
        (0, "Governed by the Connectivity Standards Alliance (CSA)"),
    ])

    add_three_columns(
        prs,
        "Zigbee Device Roles",
        [
            "Coordinator (ZC)",
            "- One per network",
            "- Forms network",
            "- Trust center function",
        ],
        [
            "Router (ZR)",
            "- Mains powered",
            "- Forwards traffic",
            "- Extends mesh coverage",
        ],
        [
            "End Device (ZED)",
            "- Battery optimized",
            "- Talks via parent",
            "- Does not route",
        ],
    )

    add_two_content(
        prs,
        "Zigbee Application Model (Product Definition Core)",
        [
            "Node -> Endpoint -> Cluster -> Attribute/Command",
            "",
            "Endpoint:",
            "- Logical function container",
            "- One product can expose multiple endpoints",
            "",
            "Cluster:",
            "- Functional module (On/Off, Level, Color, Temp)",
        ],
        [
            "Attribute:",
            "- Current state value",
            "",
            "Command:",
            "- Action request",
            "",
            "Binding & groups:",
            "- Direct control links",
            "- One-to-many room/group behavior",
            "- Core to interoperable product behavior",
        ],
    )

    add_title_content(prs, "Zigbee Working Model in Real Meshes", [
        (0, "Device joins during permit-join window and receives network parameters"),
        (0, "Gets a short NWK address for compact routing"),
        (0, "Route discovery builds multi-hop paths through routers"),
        (0, "Hop-by-hop MAC ACK supports reliability"),
        (0, "Sleepy end devices poll parent; parent buffers downlink traffic"),
        (0, "Typically connected to apps/cloud through a Zigbee gateway/hub"),
    ])

    add_title_content(prs, "Zigbee Security", [
        (0, "AES-128-CCM* integrity + encryption"),
        (0, "Network key for mesh-wide secured traffic"),
        (0, "Link keys for secure pairwise/application exchanges"),
        (0, "Trust center authorizes joins and key updates"),
        (0, "Zigbee 3.0 strengthened onboarding via install codes"),
        (0, "Legacy default-key practices are discouraged in modern deployments"),
    ])

    # 5) Thread
    add_section(prs, "Part 4: Thread", layout=6)

    add_title_content(prs, "Thread Stack and Philosophy", [
        (0, "Thread uses 802.15.4 radio/MAC plus native IPv6"),
        (0, "6LoWPAN compresses IPv6/UDP headers into small 802.15.4 frames"),
        (0, "Mesh networking with MLE and router role management"),
        (0, "No proprietary app layer required (Matter commonly used)"),
        (0, "Goal: secure, low-power, IP-native home/building mesh"),
    ])

    add_four_contents(
        prs,
        "Thread Roles",
        ["Border Router", "Bridges Thread <-> Wi-Fi/Ethernet", "Multiple can coexist for resiliency"],
        ["Leader", "Elected control-plane role", "Replaced automatically on failure"],
        ["Router / REED", "Routes traffic", "REED can promote to router as needed"],
        ["End Device (MED/SED)", "Battery role via parent", "Sleepy operation for long life"],
    )

    add_title_content(prs, "Thread Working Model: Commissioning and Data Flow", [
        (0, "Joiner device authenticates with Commissioner using onboarding secret"),
        (0, "DTLS/J-PAKE based secure commissioning exchange"),
        (0, "Device receives operational credentials and attaches to mesh"),
        (0, "Traffic is routed at IP layer across routers"),
        (0, "Border Router bridges traffic to LAN/cloud without proprietary translation"),
        (0, "This is why Thread is considered future-friendly for IP ecosystems"),
    ])

    add_title_content(prs, "Thread Security", [
        (0, "802.15.4 MAC security with AES-128-CCM for link traffic"),
        (0, "Secure commissioning before operational network access"),
        (0, "App-layer/transport security via DTLS and upper protocols"),
        (0, "Supports credential rotation and strong authenticated joining"),
        (0, "Designed with secure-by-default principles"),
    ])

    # 6) Compare, market, future
    add_section(prs, "Part 5: Compare, Matter, Market, Future", layout=6)

    add_two_content(
        prs,
        "Zigbee vs Thread: Practical Comparison",
        [
            "Zigbee",
            "- Mature ecosystem",
            "- Rich built-in app model (ZCL)",
            "- Usually requires hub/gateway",
            "- Massive installed base",
            "",
            "Best when:",
            "- Cost and proven ecosystem dominate",
        ],
        [
            "Thread",
            "- Native IPv6 mesh",
            "- Border router model",
            "- Natural fit for Matter",
            "- Strong future momentum",
            "",
            "Best when:",
            "- Future-proof multi-ecosystem interoperability is key",
        ],
    )

    add_title_content(prs, "Matter and Product Strategy", [
        (0, "Matter is the app layer over Thread/Wi-Fi"),
        (0, "BLE is used for commissioning in many Matter flows"),
        (0, "Matter data model (endpoints/clusters/attributes) maps well to Zigbee knowledge"),
        (0, "Use multiprotocol SoCs: BLE + Zigbee + Thread for portfolio flexibility"),
        (0, "Bridge legacy Zigbee installed base while shipping new Matter-over-Thread products"),
    ])

    add_title_content(prs, "Products in Market Today and Future Trends", [
        (0, "Zigbee today: major smart lighting, sensors, and hub ecosystems"),
        (0, "Thread today: growing in Apple/Google/Amazon ecosystems and Matter devices"),
        (0, "Border routers are increasingly embedded in home gateways and speakers"),
        (0, "Future: Thread + Matter growth, Zigbee continues strongly in existing deployments"),
        (0, "Convergence trend: one hardware platform, multiple protocol SKUs via firmware"),
    ])

    add_title_content(prs, "Quick Technical Checklist for Product Definition", [
        (0, "What is power class? (coin cell vs mains)"),
        (0, "Need direct phone connectivity or mesh-first architecture?"),
        (0, "Is IP-native addressing required end-to-end?"),
        (0, "Interoperability target: Zigbee ecosystem, Matter ecosystem, or both?"),
        (0, "Gateway/Border-router dependencies and failure model"),
        (0, "Security onboarding and key lifecycle requirements"),
        (0, "Regional RF/channel strategy and coexistence plan"),
    ])

    # 7) References (last slide requested)
    add_title_content(prs, "References", [
        (0, "IEEE 802.15.4-2020: https://standards.ieee.org/ieee/802.15.4/7029/"),
        (0, "CSA Zigbee overview: https://csa-iot.org/all-solutions/zigbee/"),
        (0, "Zigbee specifications portal: https://csa-iot.org/developer-resource/specifications-download-request/"),
        (0, "Thread overview: https://www.threadgroup.org/What-is-Thread/Overview"),
        (0, "Thread specs: https://www.threadgroup.org/support#specifications"),
        (0, "Matter overview: https://csa-iot.org/all-solutions/matter/"),
        (0, "OpenThread primer: https://openthread.io/guides/thread-primer"),
        (0, "RFC 4944: https://www.rfc-editor.org/rfc/rfc4944"),
        (0, "RFC 6282: https://www.rfc-editor.org/rfc/rfc6282"),
        (0, "RFC 6550: https://www.rfc-editor.org/rfc/rfc6550"),
    ])

    # 8) Closing
    add_final(prs, [
        "You are now ready to discuss and define 802.15.4 / Zigbee / Thread products.",
        "Next: Pick one target product category and derive a full feature/stack/security architecture.",
    ])

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
