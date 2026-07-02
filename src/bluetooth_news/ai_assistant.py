"""AI assistant: index local docs (PDF/DOCX/TXT/MD) + project JSON,
retrieve with BM25, fall back to web search, answer with Gemini.

Free stack:
- LLM:   Google Gemini (free tier, env GEMINI_API_KEY) -> fallback retrieval-only
- Index: rank_bm25 (no model download)
- Web:   duckduckgo-search (no key) -> trafilatura for body extraction
"""
from __future__ import annotations

import json
import logging
import os
import re
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable

log = logging.getLogger(__name__)

# Speed knobs (override via env)
_MAX_PDF_PAGES = int(os.environ.get("AI_MAX_PDF_PAGES", "80"))      # cap huge spec PDFs
_MAX_FILE_BYTES = int(os.environ.get("AI_MAX_FILE_MB", "50")) * 1024 * 1024
_INDEX_WORKERS = int(os.environ.get("AI_INDEX_WORKERS", "8"))
_SAVE_EVERY = 1000  # incremental save

# ---------------------------------------------------------------- paths
ROOT = Path(__file__).resolve().parents[2]
DATA_DIR = ROOT / "data"
OUTPUT_ROOT = ROOT / "output"
INDEX_FILE = DATA_DIR / "ai_index.json"
PREFS_FILE = DATA_DIR / "ai_prefs.json"
QA_CACHE_FILE = DATA_DIR / "ai_qa_cache.json"
RECENT_FILE = DATA_DIR / "ai_recent.json"

# Cache / recent tuning (override via env)
_QA_CACHE_MAX = int(os.environ.get("AI_QA_CACHE_MAX", "500"))
_RECENT_MAX = int(os.environ.get("AI_RECENT_MAX", "15"))

# ---------------------------------------------------------------- preferences
_DEFAULT_PREFS = {
    "primary_focus": "Bluetooth and IEEE 802.15.4 (Thread, Matter, Zigbee multiprotocol) on Infineon AIROC chips",
    "vendor": "Infineon",
    "product_family": "AIROC",
    "audience_default": "smart-home and industrial-IoT OEM decision-makers",
    "tone": "concrete, technical, marketing-savvy; cite real chips, customers, standards, dates",
    "must_avoid": ["LE Audio when the topic is XR/AR/VR",
                   "off-topic competitor product fluff",
                   "boilerplate that does not reference the actual topic"],
}


def load_prefs() -> dict:
    try:
        if PREFS_FILE.exists():
            data = json.loads(PREFS_FILE.read_text(encoding="utf-8"))
            merged = {**_DEFAULT_PREFS, **data}
            return merged
    except Exception:  # noqa: BLE001
        pass
    return dict(_DEFAULT_PREFS)


def save_prefs(updates: dict) -> dict:
    cur = load_prefs()
    cur.update({k: v for k, v in (updates or {}).items() if v is not None})
    PREFS_FILE.parent.mkdir(parents=True, exist_ok=True)
    PREFS_FILE.write_text(json.dumps(cur, indent=2), encoding="utf-8")
    return cur


def prefs_block(label: str = "USER PREFERENCES") -> str:
    """Render the prefs as a system-prompt block."""
    p = load_prefs()
    return (
        f"\n\n{label} (always honor these unless the user explicitly overrides):\n"
        f"- Primary focus: {p['primary_focus']}.\n"
        f"- Speaking on behalf of vendor: {p['vendor']} (product family: {p['product_family']}).\n"
        f"- Default audience: {p['audience_default']}.\n"
        f"- Tone: {p['tone']}.\n"
        f"- Must avoid: " + "; ".join(p.get("must_avoid", [])) + ".\n"
    )


# ---------------------------------------------------------------- QA cache
_QA_STOPWORDS = {
    "the", "a", "an", "of", "for", "to", "and", "or", "is", "are", "was", "were",
    "in", "on", "at", "by", "with", "me", "my", "our", "us", "you", "your",
    "what", "whats", "how", "do", "does", "did", "can", "could", "should",
    "please", "tell", "give", "show", "about", "vs",
}


def _qa_fingerprint(question: str) -> str:
    """Normalized bag-of-words key so trivial rephrasings hit the same cache slot."""
    toks = re.findall(r"[a-z0-9][a-z0-9\.\-]*", (question or "").lower())
    toks = [t for t in toks if t not in _QA_STOPWORDS]
    return " ".join(sorted(set(toks)))


def _index_signature() -> str:
    """Changes whenever the index is rebuilt — used to invalidate stale answers."""
    try:
        return str(int(INDEX_FILE.stat().st_mtime))
    except OSError:
        return "0"


def qa_cache_get(question: str) -> dict | None:
    fp = _qa_fingerprint(question)
    if not fp or not QA_CACHE_FILE.exists():
        return None
    try:
        cache = json.loads(QA_CACHE_FILE.read_text(encoding="utf-8"))
    except Exception:
        return None
    entry = cache.get(fp)
    if not entry:
        return None
    if entry.get("sig") != _index_signature():
        return None  # index changed since answer was cached
    return entry


def qa_cache_put(question: str, answer: "Answer") -> None:
    fp = _qa_fingerprint(question)
    if not fp:
        return
    cache: dict = {}
    if QA_CACHE_FILE.exists():
        try:
            cache = json.loads(QA_CACHE_FILE.read_text(encoding="utf-8"))
        except Exception:
            cache = {}
    cache[fp] = {
        "question": question,
        "reply": answer.reply,
        "sources": answer.sources,
        "backend": answer.backend,
        "used_web": answer.used_web,
        "sig": _index_signature(),
        "ts": time.time(),
    }
    # LRU cap: drop oldest by timestamp.
    if len(cache) > _QA_CACHE_MAX:
        for k in sorted(cache, key=lambda k: cache[k].get("ts", 0))[: len(cache) - _QA_CACHE_MAX]:
            cache.pop(k, None)
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    try:
        QA_CACHE_FILE.write_text(json.dumps(cache, ensure_ascii=False), encoding="utf-8")
    except Exception as e:  # noqa: BLE001
        log.warning("qa cache write failed: %s", e)


# ---------------------------------------------------------------- recent questions
def load_recent(limit: int = _RECENT_MAX) -> list[str]:
    if not RECENT_FILE.exists():
        return []
    try:
        data = json.loads(RECENT_FILE.read_text(encoding="utf-8"))
        items = data.get("recent", []) if isinstance(data, dict) else list(data)
        return [str(x) for x in items][:limit]
    except Exception:
        return []


def push_recent(question: str, limit: int = _RECENT_MAX) -> list[str]:
    q = (question or "").strip()
    if not q:
        return load_recent(limit)
    items = load_recent(200)
    # de-dupe case-insensitively, newest first
    items = [x for x in items if x.strip().lower() != q.lower()]
    items.insert(0, q)
    items = items[:limit]
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    try:
        RECENT_FILE.write_text(json.dumps({"recent": items}, ensure_ascii=False),
                               encoding="utf-8")
    except Exception as e:  # noqa: BLE001
        log.warning("recent write failed: %s", e)
    return items


def clear_recent() -> None:
    try:
        RECENT_FILE.write_text(json.dumps({"recent": []}), encoding="utf-8")
    except Exception:
        pass


# ---------------------------------------------------------------- suggested questions
_CURATED_SUGGESTIONS = [
    "Teach me IEEE 802.15.4 — protocol stack, key features, and how it differs from Bluetooth LE.",
    "Summarise Bluetooth 6.0 Channel Sounding and our positioning vs Nordic and Silicon Labs.",
    "What are the latest news highlights across our tracked IoT-wireless topics?",
    "Give me a Matter 1.4 + Thread 1.4 readiness overview for smart-home OEMs.",
]


def dynamic_suggestions(limit: int = 8) -> list[str]:
    """Build starter questions from indexed data so they stay relevant, blended
    with a curated base set."""
    out: list[str] = list(_CURATED_SUGGESTIONS)
    # From customers.json -> opportunity questions.
    try:
        cust_p = DATA_DIR / "customers.json"
        if cust_p.exists():
            cust = json.loads(cust_p.read_text(encoding="utf-8")).get("customers", [])
            for c in cust[:4]:
                name = c.get("name")
                if name:
                    out.append(f"Which wireless chip vendors does {name} use, and where could AIROC win?")
    except Exception:
        pass
    # From competitors.json -> comparison questions.
    try:
        comp_p = DATA_DIR / "competitors.json"
        if comp_p.exists():
            comp = json.loads(comp_p.read_text(encoding="utf-8")).get("competitors", [])
            for c in comp[:4]:
                v = c.get("vendor")
                if v:
                    out.append(f"How does Infineon AIROC compare to {v}'s latest wireless SoCs?")
    except Exception:
        pass
    # de-dupe preserving order
    seen: set[str] = set()
    uniq = [q for q in out if not (q in seen or seen.add(q))]
    return uniq[:limit]


# Always-scanned doc roots (in addition to AI_EXTRA_DOC_PATHS). These hold
# Bluetooth SIG specifications & related technical material.
_DEFAULT_DOC_ROOTS = [
    Path(r"C:\guptakanak\SIG"),
]


# Allow user to add more roots via env var (semicolon-separated on Windows).
def _extra_doc_roots() -> list[Path]:
    out: list[Path] = list(_DEFAULT_DOC_ROOTS)
    raw = os.environ.get("AI_EXTRA_DOC_PATHS", "")
    for p in re.split(r"[;,\n]+", raw):
        p = p.strip().strip('"')
        if p:
            out.append(Path(p))
    # de-dupe (case-insensitive on Windows) while preserving order
    seen: set[str] = set()
    res: list[Path] = []
    for p in out:
        k = str(p).lower()
        if k not in seen:
            seen.add(k)
            res.append(p)
    return res


# ---------------------------------------------------------------- chunking
@dataclass
class Chunk:
    text: str
    source: str          # file path or "data:customers" or "web:<url>"
    title: str = ""
    page: int = 0
    kind: str = "doc"    # doc | data | web | news

    def to_dict(self) -> dict:
        return self.__dict__


_WORD_RE = re.compile(r"\S+")


def _split(text: str, max_words: int = 380, overlap: int = 60) -> list[str]:
    words = _WORD_RE.findall(text)
    if not words:
        return []
    out = []
    i = 0
    n = len(words)
    while i < n:
        out.append(" ".join(words[i:i + max_words]))
        if i + max_words >= n:
            break
        i += max_words - overlap
    return out


# ---------------------------------------------------------------- file readers
def _read_pdf(path: Path) -> Iterable[tuple[int, str]]:
    try:
        from pypdf import PdfReader
    except ImportError:
        log.warning("pypdf not installed; skipping %s", path.name)
        return
    try:
        r = PdfReader(str(path))
        for i, page in enumerate(r.pages, 1):
            if i > _MAX_PDF_PAGES:
                break
            try:
                t = page.extract_text() or ""
            except Exception:
                t = ""
            if t.strip():
                yield i, t
    except Exception as e:
        log.warning("PDF read failed %s: %s", path.name, e)


def _read_docx(path: Path) -> Iterable[tuple[int, str]]:
    try:
        from docx import Document
    except ImportError:
        log.warning("python-docx not installed; skipping %s", path.name)
        return
    try:
        d = Document(str(path))
        parts = [p.text for p in d.paragraphs if p.text]
        for tbl in d.tables:
            for row in tbl.rows:
                parts.append(" | ".join(c.text for c in row.cells))
        t = "\n".join(parts)
        if t.strip():
            yield 1, t
    except Exception as e:
        log.warning("DOCX read failed %s: %s", path.name, e)


def _read_doc_legacy(_path: Path) -> Iterable[tuple[int, str]]:
    # Legacy binary .doc is not supported by python-docx; skip silently.
    return
    yield  # pragma: no cover


def _read_text(path: Path) -> Iterable[tuple[int, str]]:
    try:
        t = path.read_text(encoding="utf-8", errors="ignore")
        if t.strip():
            yield 1, t
    except Exception as e:
        log.warning("Text read failed %s: %s", path.name, e)


def _read_pptx(path: Path) -> Iterable[tuple[int, str]]:
    try:
        from pptx import Presentation
    except ImportError:
        log.warning("python-pptx not installed; skipping %s", path.name)
        return
    try:
        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides, 1):
            parts: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs)
                        if t.strip():
                            parts.append(t)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        parts.append(" | ".join(c.text for c in row.cells))
            txt = "\n".join(parts)
            if txt.strip():
                yield i, txt
    except Exception as e:
        log.warning("PPTX read failed %s: %s", path.name, e)


_READERS = {
    ".pdf": _read_pdf,
    ".docx": _read_docx,
    ".doc": _read_doc_legacy,
    ".pptx": _read_pptx,
    ".txt": _read_text,
    ".md": _read_text,
    ".markdown": _read_text,
    ".json": _read_text,
    ".csv": _read_text,
}


# ---------------------------------------------------------------- version-aware dedup
_VER_RE = re.compile(r"v(\d+(?:\.\d+)+)", re.I)   # v5.4, v6.1, v1.0 ...
_COPY_RE = re.compile(r"\s*\(\d+\)\s*$")           # trailing " (1)", " (002)"


def _norm_words(s: str) -> str:
    return re.sub(r"[ _\-]+", " ", s).strip(" _-.").lower()


def _dedup_group_key(p: Path):
    """Files that share this key are treated as versions/copies of the same doc.
    Versioned files (Core_v5.4.pdf) group across folders; unversioned files only
    collapse exact copies within the same folder."""
    stem = _COPY_RE.sub("", p.stem)
    m = _VER_RE.search(stem)
    if m:
        base = _norm_words(_VER_RE.sub("", stem))
        return ("v", base, p.suffix.lower())
    return ("p", str(p.parent).lower(), _norm_words(stem), p.suffix.lower())


def _version_tuple(p: Path) -> tuple:
    m = _VER_RE.search(p.stem)
    return tuple(int(x) for x in m.group(1).split(".")) if m else ()


def _dedup_score(p: Path) -> tuple:
    is_copy = 1 if _COPY_RE.search(p.stem) else 0
    try:
        mtime = p.stat().st_mtime
    except OSError:
        mtime = 0.0
    # highest version wins; prefer non-copy; then newest.
    return (_version_tuple(p), -is_copy, mtime)


def _dedup_versions(files: list[Path]) -> list[Path]:
    best: dict = {}
    for p in files:
        k = _dedup_group_key(p)
        if k not in best or _dedup_score(p) > _dedup_score(best[k]):
            best[k] = p
    dropped = len(files) - len(best)
    if dropped:
        log.info("version dedup: kept %d files, dropped %d older/duplicate versions",
                 len(best), dropped)
    return list(best.values())


# ---------------------------------------------------------------- index build
def _collect_files(roots: list[Path]) -> list[Path]:
    seen: set[Path] = set()
    files: list[Path] = []
    for root in roots:
        if not root.exists():
            log.info("Doc root missing, skipping: %s", root)
            continue
        for p in root.rglob("*"):
            if not p.is_file():
                continue
            if p.name.startswith("~$"):  # Office lock files
                continue
            try:
                rp = p.resolve()
            except OSError:
                continue
            if rp in seen:
                continue
            ext = p.suffix.lower()
            if ext not in _READERS:
                continue
            try:
                if p.stat().st_size > _MAX_FILE_BYTES:
                    continue
            except OSError:
                continue
            seen.add(rp)
            files.append(p)
    return _dedup_versions(files)


def _read_one(p: Path) -> list[Chunk]:
    reader = _READERS.get(p.suffix.lower())
    if not reader:
        return []
    out: list[Chunk] = []
    rel = str(p)
    for page, text in reader(p):
        for chunk_text in _split(text):
            out.append(Chunk(text=chunk_text, source=rel, title=p.name,
                             page=page, kind="doc"))
    return out


def _ingest_project_json() -> Iterable[Chunk]:
    """Pull rich text from data/*.json so the assistant 'knows' our customers/competitors."""
    for fname in ("customers.json", "competitors.json", "events.json",
                  "standards.json", "external.json", "github_repos.json",
                  "briefing_input.json", "roadmap_input.json"):
        p = DATA_DIR / fname
        if not p.exists():
            continue
        try:
            obj = json.loads(p.read_text(encoding="utf-8"))
        except Exception:
            continue
        flat = json.dumps(obj, ensure_ascii=False, indent=2)
        for chunk_text in _split(flat, max_words=500, overlap=50):
            yield Chunk(text=chunk_text, source=f"data:{fname}",
                        title=fname, kind="data")


# ---------------------------------------------------------------- report HTML ingest
def _latest_report_dir() -> Path | None:
    """The report bundle the AI server actually serves (newest site_* dir),
    falling back to output/latest."""
    if not OUTPUT_ROOT.exists():
        return None
    sites = sorted(
        (p for p in OUTPUT_ROOT.glob("site_*") if p.is_dir()),
        key=lambda p: p.stat().st_mtime, reverse=True,
    )
    if sites:
        return sites[0]
    latest = OUTPUT_ROOT / "latest"
    return latest if latest.is_dir() else None


_HTML_TAG_RE = re.compile(r"<[^>]+>")


def _html_to_text(html_src: str) -> str:
    """Strip a report HTML page down to readable text, preserving table structure."""
    import html as _htmlmod
    s = html_src
    # Drop non-content blocks entirely.
    s = re.sub(r"(?is)<script.*?</script>", " ", s)
    s = re.sub(r"(?is)<style.*?</style>", " ", s)
    s = re.sub(r"(?is)<head.*?</head>", " ", s)
    s = re.sub(r"(?is)<nav.*?</nav>", " ", s)
    s = re.sub(r"(?is)<svg.*?</svg>", " ", s)
    # Keep table cell/row structure so figures stay associated with labels.
    s = re.sub(r"(?is)</(td|th)>", " | ", s)
    s = re.sub(r"(?is)</(tr|table|div|section|article|p|li|h[1-6])>", "\n", s)
    s = _HTML_TAG_RE.sub(" ", s)
    s = _htmlmod.unescape(s)
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n\s*\n+", "\n", s)
    return s.strip()


def _ingest_report_html() -> Iterable[Chunk]:
    """Index the rendered report pages (output/.../*.html) so the assistant sees
    exactly what a human sees on the site — customers, competitors, news, roadmap, etc."""
    d = _latest_report_dir()
    if not d:
        return
    for p in sorted(d.glob("*.html")):
        if p.name.lower() == "ai.html":  # skip the chat shim itself
            continue
        try:
            raw = p.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            continue
        text = _html_to_text(raw)
        if not text or len(text) < 40:
            continue
        title = p.stem.replace("_", " ").replace("-", " ").title()
        for chunk_text in _split(text, max_words=450, overlap=60):
            yield Chunk(text=chunk_text, source=f"report:{p.name}",
                        title=title, kind="report")


def _save(chunks: list[Chunk]) -> None:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    INDEX_FILE.write_text(
        json.dumps({"chunks": [c.to_dict() for c in chunks]}, ensure_ascii=False),
        encoding="utf-8",
    )


def build_index(verbose: bool = True) -> dict:
    roots = _extra_doc_roots()
    if verbose:
        print(f"[ai] scanning roots: {[str(r) for r in roots]}")
    files = _collect_files(roots)
    if verbose:
        print(f"[ai] {len(files)} files to parse — {_INDEX_WORKERS} workers, "
              f"PDF page cap={_MAX_PDF_PAGES}, max size={_MAX_FILE_BYTES//1024//1024}MB")

    chunks: list[Chunk] = []
    done = 0
    with ThreadPoolExecutor(max_workers=_INDEX_WORKERS) as ex:
        futures = {ex.submit(_read_one, p): p for p in files}
        for fut in as_completed(futures):
            try:
                chunks.extend(fut.result())
            except Exception as e:
                log.warning("reader crashed on %s: %s", futures[fut].name, e)
            done += 1
            if verbose and done % 25 == 0:
                print(f"[ai]   files {done}/{len(files)}  chunks={len(chunks)}")
            if len(chunks) and len(chunks) % _SAVE_EVERY < 50:
                _save(chunks)  # incremental snapshot

    if verbose:
        print(f"[ai] doc chunks: {len(chunks)}")
    pre = len(chunks)
    for c in _ingest_project_json():
        chunks.append(c)
    if verbose:
        print(f"[ai] data chunks: {len(chunks) - pre}")
    pre_html = len(chunks)
    for c in _ingest_report_html():
        chunks.append(c)
    if verbose:
        print(f"[ai] report chunks: {len(chunks) - pre_html}")
    _save(chunks)
    if verbose:
        print(f"[ai] wrote {INDEX_FILE} ({len(chunks)} chunks)")
    return {"count": len(chunks), "path": str(INDEX_FILE)}


# ---------------------------------------------------------------- retrieval
class Retriever:
    def __init__(self) -> None:
        self.chunks: list[Chunk] = []
        self.bm25 = None
        self._tokens: list[list[str]] = []
        self.load()

    def load(self) -> None:
        if not INDEX_FILE.exists():
            log.info("No AI index yet; run build_index().")
            return
        try:
            obj = json.loads(INDEX_FILE.read_text(encoding="utf-8"))
        except Exception as e:
            log.warning("Index load failed: %s", e)
            return
        self.chunks = [Chunk(**c) for c in obj.get("chunks", [])]
        self._fit()

    def _fit(self) -> None:
        if not self.chunks:
            self.bm25 = None
            return
        try:
            from rank_bm25 import BM25Okapi
        except ImportError:
            log.warning("rank_bm25 not installed; falling back to substring search")
            self.bm25 = None
            return
        self._tokens = [self._tokenize(c.text) for c in self.chunks]
        self.bm25 = BM25Okapi(self._tokens)

    @staticmethod
    def _tokenize(text: str) -> list[str]:
        return re.findall(r"[A-Za-z0-9][A-Za-z0-9\-\.]+", text.lower())

    def search(self, query: str, k: int = 6) -> list[tuple[Chunk, float]]:
        if not self.chunks:
            return []
        q = self._tokenize(query)
        if not q:
            return []
        if self.bm25 is not None:
            scores = self.bm25.get_scores(q)
            ranked = sorted(enumerate(scores), key=lambda x: x[1], reverse=True)[:k]
            return [(self.chunks[i], float(s)) for i, s in ranked if s > 0]
        # naive fallback
        scored = []
        ql = set(q)
        for i, toks in enumerate(self._tokens or [self._tokenize(c.text) for c in self.chunks]):
            sc = sum(1 for t in toks if t in ql)
            if sc:
                scored.append((i, sc))
        scored.sort(key=lambda x: x[1], reverse=True)
        return [(self.chunks[i], float(s)) for i, s in scored[:k]]


# ---------------------------------------------------------------- web search
def web_search(query: str, max_results: int = 4) -> list[Chunk]:
    try:
        from ddgs import DDGS  # current package name
    except ImportError:
        try:
            from duckduckgo_search import DDGS  # legacy fallback
        except ImportError:
            log.info("ddgs / duckduckgo-search not installed; web disabled")
            return []
    results: list[Chunk] = []
    try:
        with DDGS() as ddgs:
            hits = list(ddgs.text(query, max_results=max_results, region="wt-wt"))
    except Exception as e:
        log.warning("Web search failed: %s", e)
        return []
    try:
        import trafilatura
    except ImportError:
        trafilatura = None  # type: ignore
    for h in hits:
        url = h.get("href") or h.get("url") or ""
        title = h.get("title") or url
        body = h.get("body") or ""
        if trafilatura and url:
            try:
                downloaded = trafilatura.fetch_url(url, no_ssl=True)
                if downloaded:
                    extracted = trafilatura.extract(downloaded, include_comments=False,
                                                    include_tables=False, favor_recall=True)
                    if extracted:
                        body = extracted[:4000]
            except Exception:
                pass
        if not body:
            continue
        for chunk_text in _split(body, max_words=400, overlap=40):
            results.append(Chunk(text=chunk_text, source=f"web:{url}",
                                 title=title, kind="web"))
    return results


def fetch_url(url: str, max_chars: int = 8000) -> list[Chunk]:
    """Directly fetch & extract a single URL. Returns Chunk(s) ready to ingest."""
    if not url:
        return []
    try:
        import trafilatura
    except ImportError:
        return []
    try:
        downloaded = trafilatura.fetch_url(url, no_ssl=True)
        if not downloaded:
            return []
        title = url
        try:
            meta = trafilatura.extract_metadata(downloaded)
            if meta and getattr(meta, "title", None):
                title = meta.title
        except Exception:
            pass
        text = trafilatura.extract(downloaded, include_comments=False,
                                   include_tables=True, favor_recall=True)
        if not text:
            return []
        text = text[:max_chars]
        return [Chunk(text=t, source=f"web:{url}", title=title, kind="web")
                for t in _split(text, max_words=400, overlap=40)]
    except Exception as e:
        log.warning("fetch_url failed (%s): %s", url[:80], e)
        return []


def web_image_search(query: str, max_results: int = 5,
                     min_width: int = 600, min_height: int = 400) -> list[dict]:
    """Free image search via DuckDuckGo + Wikimedia Commons fallback.
    Returns list of {url, source_page, title, width, height, source}. No keys needed."""
    out: list[dict] = []
    # DuckDuckGo Images
    try:
        try:
            from ddgs import DDGS  # newer package name
        except ImportError:
            from duckduckgo_search import DDGS
        with DDGS() as ddgs:
            hits = list(ddgs.images(query, max_results=max_results * 2,
                                    safesearch="moderate", size="Large",
                                    type_image="photo", license_image="any"))
        for h in hits:
            url = h.get("image") or h.get("thumbnail") or ""
            w = int(h.get("width") or 0)
            ht = int(h.get("height") or 0)
            if not url or w < min_width or ht < min_height:
                continue
            out.append({"url": url, "source_page": h.get("url", ""),
                        "title": h.get("title", "")[:120],
                        "width": w, "height": ht, "source": "ddg"})
            if len(out) >= max_results:
                break
    except Exception as e:
        log.warning("ddg image search failed: %s", e)

    # Wikimedia Commons fallback (always tried — gives clean, royalty-free imagery)
    if len(out) < max_results:
        try:
            import requests as _rq
            r = _rq.get("https://commons.wikimedia.org/w/api.php",
                        params={"action": "query", "format": "json",
                                "generator": "search", "gsrsearch": f"filetype:bitmap {query}",
                                "gsrlimit": max_results * 2,
                                "prop": "imageinfo", "iiprop": "url|size",
                                "iiurlwidth": 1280},
                        headers={"User-Agent": "AIROC-AI/1.0 (research assistant)"},
                        timeout=8)
            data = r.json().get("query", {}).get("pages", {}) if r.ok else {}
            for _pid, page in data.items():
                ii = (page.get("imageinfo") or [{}])[0]
                url = ii.get("thumburl") or ii.get("url")
                w = int(ii.get("thumbwidth") or ii.get("width") or 0)
                ht = int(ii.get("thumbheight") or ii.get("height") or 0)
                if not url or w < min_width:
                    continue
                out.append({"url": url, "source_page": "https://commons.wikimedia.org/wiki/" + page.get("title", "").replace(" ", "_"),
                            "title": page.get("title", "")[:120],
                            "width": w, "height": ht, "source": "wikimedia"})
                if len(out) >= max_results:
                    break
        except Exception as e:
            log.warning("wikimedia image search failed: %s", e)
    return out


def ingest_web_results(chunks: list[Chunk]) -> int:
    """Append web chunks to the saved index so future queries reuse them."""
    if not chunks:
        return 0
    obj = {"chunks": []}
    if INDEX_FILE.exists():
        try:
            obj = json.loads(INDEX_FILE.read_text(encoding="utf-8"))
        except Exception:
            pass
    seen = {(c.get("source"), c.get("text")[:80]) for c in obj.get("chunks", [])}
    added = 0
    for c in chunks:
        key = (c.source, c.text[:80])
        if key in seen:
            continue
        obj.setdefault("chunks", []).append(c.to_dict())
        seen.add(key)
        added += 1
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    INDEX_FILE.write_text(json.dumps(obj, ensure_ascii=False), encoding="utf-8")
    return added


# ---------------------------------------------------------------- web enrichment (nightly)
_DEFAULT_SEED_QUERIES = [
    # Standards / tech
    "Bluetooth 6.0 Channel Sounding latest news",
    "Bluetooth LE Audio Auracast 2026 deployments",
    "IEEE 802.15.4 multiprotocol SoC roadmap",
    "Thread 1.4 border router updates",
    "Matter 1.4 device certifications",
    "Aliro digital key NFC UWB BLE",
    "Wi-Fi 7 MLO chipset launches 2026",
    "Wi-Fi 8 802.11bn timeline",
    # Infineon
    "Infineon AIROC CYW55513 design wins",
    "Infineon AIROC Bluetooth 6.0 announcements",
    "Infineon CYW20829 Thread Matter",
    # Competitors
    "NXP IW612 IW610 design wins 2026",
    "Qualcomm QCC74x Wi-Fi BT smart home",
    "Silicon Labs SiWx917 Matter",
    "Nordic nRF54 series news",
    "Espressif ESP32-C6 ESP32-H2 Matter Thread",
    "STMicroelectronics STM32WB ST33 Bluetooth",
    "Synaptics SYN43 Veros wireless",
    "Renesas DA147xx wireless MCU",
    "MediaTek Filogic Wi-Fi 7",
    "TI CC2340 CC2674 Bluetooth Matter",
    # Customers (no auto)
    "Google Nest Hub Matter Thread chipset",
    "Amazon Echo new device 2026 wireless chip",
    "Meta Quest Ray-Ban smart glasses connectivity",
    "Arlo camera new model 2026 Wi-Fi chipset",
    "Stryker medical wearable Bluetooth",
    "Motorola wearable Bluetooth 2026",
    "Bose Sonos Ultimate Ears LE Audio Auracast",
    "Garmin Whoop Oura wearable wireless chipset",
    "Samsung Galaxy Watch Bluetooth LE Audio",
    "Apple AirPods Vision Pro wireless chipset",
]


def enrich_from_web(queries: list[str] | None = None,
                    per_query: int = 3,
                    include_data_seeds: bool = True,
                    verbose: bool = True) -> dict:
    """Pull a fresh batch of web snippets for each seed query and append to index.

    Used by the nightly scheduled task. Safe to call repeatedly - duplicates skipped.
    When include_data_seeds=True, also auto-generates queries from
    data/customers.json + data/competitors.json so additions there flow through.
    """
    queries = list(queries) if queries else list(_DEFAULT_SEED_QUERIES)
    if include_data_seeds:
        queries.extend(_data_driven_seed_queries())
    # de-dup while preserving order
    seen: set[str] = set()
    queries = [q for q in queries if not (q in seen or seen.add(q))]
    total_added = 0
    total_chunks = 0
    failures = 0
    for i, q in enumerate(queries, 1):
        try:
            ch = web_search(q, max_results=per_query)
        except Exception as e:
            failures += 1
            if verbose:
                print(f"[ai-enrich] [{i}/{len(queries)}] FAILED '{q}': {e}")
            continue
        added = ingest_web_results(ch)
        total_chunks += len(ch)
        total_added += added
        if verbose:
            print(f"[ai-enrich] [{i}/{len(queries)}] '{q[:60]}' -> "
                  f"{len(ch)} chunks, {added} new")
    return {"queries": len(queries), "chunks_fetched": total_chunks,
            "chunks_added": total_added, "failures": failures}


def _data_driven_seed_queries() -> list[str]:
    """Build extra queries from data/customers.json + data/competitors.json."""
    out: list[str] = []
    cust_p = DATA_DIR / "customers.json"
    if cust_p.exists():
        try:
            cust = json.loads(cust_p.read_text(encoding="utf-8")).get("customers", [])
            for c in cust:
                name = c.get("name")
                if name:
                    out.append(f"{name} 2026 wireless chipset Bluetooth Wi-Fi")
                    out.append(f"{name} new product launch IoT connectivity")
        except Exception:
            pass
    comp_p = DATA_DIR / "competitors.json"
    if comp_p.exists():
        try:
            comp = json.loads(comp_p.read_text(encoding="utf-8"))
            for c in comp.get("competitors", []):
                v = c.get("vendor")
                if v:
                    out.append(f"{v} Bluetooth Wi-Fi 802.15.4 design wins 2026")
                    out.append(f"{v} new SoC announcement smart home")
        except Exception:
            pass
    return out


# ---------------------------------------------------------------- LLM (Gemini)
SYSTEM_PROMPT = """You are 'AIROC AI' — a senior IoT-wireless marketing & product strategist
working for Infineon. You help with Bluetooth, Wi-Fi, IEEE 802.15.4, Thread, Matter,
Aliro, UWB. You know Infineon's AIROC family and major competitors (NXP, STMicro,
Qualcomm, Synaptics, Renesas, TI, Silicon Labs, Nordic, Espressif, MediaTek, Realtek).
You know top customers (Google, Amazon, Apple, Meta, Microsoft, Samsung, Arlo,
Stryker, Motorola, Bose, Sonos, Garmin, etc.) and their wireless product lines.

IMPORTANT — you DO have access to live web content. The runtime that calls you
performs live DuckDuckGo searches and fetches vendor / standards-body pages
(bluetooth.com, ieee.org, csa-iot.org, threadgroup.org, wi-fi.org, infineon.com,
nxp.com, silabs.com, nordicsemi.com, qualcomm.com, ti.com, renesas.com,
espressif.com, st.com, etc.) and inserts the extracted text into your CONTEXT
as items marked '(web)'. NEVER tell the user you cannot browse the internet or
access vendor websites — instead, USE the web context items in CONTEXT, cite
them as [#n], and if a specific page is missing, say "the runtime did not fetch
that specific page — ask me to look up X on Y.com" so I can re-run the search.

Rules:
1. Ground every factual claim in the provided CONTEXT first. Cite sources inline as
   [#1], [#2] matching the numbered context items.
2. If CONTEXT is thin, say so and use general knowledge — but flag those parts as
   "(general knowledge, not from indexed sources)".
3. When asked to teach a spec or standard (e.g. Bluetooth HDT, Channel Sounding,
   LE Audio, Auracast, 802.15.4, Thread, Matter, Wi-Fi 7), be EXHAUSTIVE on
   numerical/technical details. ALWAYS enumerate, when applicable:
     - All PHY layers and data rates (e.g. 1 Mbps, 2 Mbps, 3 Mbps, 4 Mbps, 8 Mbps).
     - Modulation schemes (e.g. GFSK, π/4-DQPSK, 8DPSK, OFDM).
     - Frequency bands and channel bandwidth.
     - Packet formats / max payload sizes.
     - Spec version introduced and backwards compatibility.
     - Power class / TX power range.
   Present this as a compact comparison table when there are multiple options.
4. Structure teaching answers: TL;DR → Key concepts → Technical details (table) →
   Spec versions → Real-world use → Infineon angle → Competitive notes.
5. When asked for market strategy output, provide concise and actionable sections:
    Situation, Key Signals, Recommendation, Risks, Next Actions.
6. Be concise, use bullet lists and short paragraphs, prefer tables for comparisons.
7. NEVER invent numbers. If you're not confident about a specific data rate or
   modulation, say "(verify on bluetooth.com / IEEE spec)" rather than guessing.
"""


# Free-tier-friendly Gemini fallback chain. Order = try first, fall back on 429/error.
# Verified May 2026: legacy 1.5 models are retired (404). Gemma models share a
# different quota pool from Gemini, so they're a useful last resort when Gemini
# free-tier is exhausted for the day.
_DEFAULT_GEMINI_CHAIN = [
    "gemini-2.5-flash-lite",     # generous free quota
    "gemini-2.0-flash-lite",
    "gemini-flash-lite-latest",
    "gemini-2.5-flash",
    "gemini-2.0-flash",
    "gemini-flash-latest",
    "gemma-3-27b-it",            # separate quota pool
    "gemma-3-12b-it",
]


def _gemini_models():
    """Return an ordered list of GenerativeModel instances to try."""
    api_key = os.environ.get("GEMINI_API_KEY") or os.environ.get("GOOGLE_API_KEY")
    if not api_key:
        return []
    try:
        import google.generativeai as genai
    except ImportError:
        log.warning("google-generativeai not installed; LLM disabled")
        return []
    genai.configure(api_key=api_key)
    # Allow override: GEMINI_MODELS="m1,m2,m3" or single GEMINI_MODEL=m
    raw = os.environ.get("GEMINI_MODELS") or os.environ.get("GEMINI_MODEL")
    if raw:
        names = [n.strip() for n in raw.split(",") if n.strip()]
        # Always append the default chain after user picks (dedup, preserve order)
        for n in _DEFAULT_GEMINI_CHAIN:
            if n not in names:
                names.append(n)
    else:
        names = list(_DEFAULT_GEMINI_CHAIN)
    out = []
    for n in names:
        # Gemma models on the Gemini API do not support system_instruction.
        if n.lower().startswith("gemma"):
            out.append(genai.GenerativeModel(n))
        else:
            out.append(genai.GenerativeModel(n, system_instruction=SYSTEM_PROMPT + prefs_block()))
    return out


def _gemini_model():
    """Back-compat: return first model in the chain (or None)."""
    models = _gemini_models()
    return models[0] if models else None


def _format_context(items: list[tuple[Chunk, float]]) -> tuple[str, list[dict]]:
    lines = []
    cites = []
    for i, (c, s) in enumerate(items, 1):
        loc = f" p.{c.page}" if c.kind == "doc" and c.page else ""
        head = f"[#{i}] {c.title or c.source}{loc} ({c.kind})"
        lines.append(f"{head}\n{c.text}\n")
        cites.append({"id": i, "title": c.title or c.source, "source": c.source,
                      "kind": c.kind, "page": c.page, "score": round(s, 2)})
    return "\n---\n".join(lines), cites


@dataclass
class Answer:
    reply: str
    sources: list[dict] = field(default_factory=list)
    used_web: bool = False
    backend: str = "retrieval-only"


# --- spec/teaching deep-research helpers --------------------------------
_SPEC_KEYWORDS = (
    "hdt", "high data throughput", "channel sounding", "le audio", "auracast",
    "isochronous", "iso channel", "leaudio", "broadcast audio", "lc3",
    "bluetooth 5", "bluetooth 6", "bt 5", "bt 6", "bt6", "bt5",
    "802.15.4", "ieee 802", "thread 1.", "matter 1.", "aliro",
    "wi-fi 6", "wi-fi 7", "wi-fi 8", "wifi 7", "wifi 8", "802.11",
    "phy", "modulation", "data rate", "mbps", "throughput", "spec",
    "uwb", "ultra wideband",
)
_TEACH_TRIGGERS = ("teach", "explain", "what is", "what's", "how does", "describe",
                   "walk me through", "deep dive", "overview of", "tell me about")
_AUTHORITATIVE_DOMAINS = (
    "bluetooth.com", "ieee.org", "csa-iot.org", "threadgroup.org",
    "wi-fi.org", "wifialliance.org", "ietf.org", "fira.org",
    "infineon.com", "nxp.com", "silabs.com", "nordicsemi.com",
    "qualcomm.com", "ti.com", "espressif.com", "renesas.com",
    "st.com", "mediatek.com", "realtek.com", "synaptics.com",
)

# Vendor name → official site, used for site:-scoped queries when the user asks
# about a specific competitor's product / portfolio.
_VENDOR_DOMAINS = {
    "infineon": "infineon.com",
    "airoc": "infineon.com",
    "nxp": "nxp.com",
    "silicon labs": "silabs.com",
    "silabs": "silabs.com",
    "nordic": "nordicsemi.com",
    "qualcomm": "qualcomm.com",
    "qca": "qualcomm.com",
    "ti": "ti.com",
    "texas instruments": "ti.com",
    "espressif": "espressif.com",
    "renesas": "renesas.com",
    "stmicro": "st.com",
    "st microelectronics": "st.com",
    "mediatek": "mediatek.com",
    "realtek": "realtek.com",
    "synaptics": "synaptics.com",
}


def _vendor_in_question(q: str) -> list[str]:
    s = (q or "").lower()
    hits = []
    for k, dom in _VENDOR_DOMAINS.items():
        if re.search(r"\b" + re.escape(k) + r"\b", s):
            if dom not in hits:
                hits.append(dom)
    return hits


def _is_spec_question(q: str) -> bool:
    s = (q or "").lower()
    if any(t in s for t in _TEACH_TRIGGERS):
        return True
    if _vendor_in_question(q):
        return True
    return any(k in s for k in _SPEC_KEYWORDS)


def _deep_web_research(question: str, max_per_query: int = 3) -> list[Chunk]:
    """For spec/teaching questions, fan out targeted queries to authoritative sources."""
    q = question.strip()
    qs = [
        q,
        f"{q} site:bluetooth.com" if "bluetooth" in q.lower() or "ble" in q.lower() or "hdt" in q.lower() or "auracast" in q.lower() or "le audio" in q.lower() else f"{q} specification",
        f"{q} PHY data rate Mbps modulation",
        f"{q} specification overview 2025",
    ]
    if "15.4" in q.lower() or "thread" in q.lower() or "zigbee" in q.lower():
        qs.append(f"{q} site:threadgroup.org OR site:csa-iot.org")
    if "matter" in q.lower():
        qs.append(f"{q} site:csa-iot.org")
    if "wi-fi" in q.lower() or "wifi" in q.lower() or "802.11" in q.lower():
        qs.append(f"{q} site:wi-fi.org")
    if "uwb" in q.lower() or "ultra wideband" in q.lower():
        qs.append(f"{q} site:fira.org")
    # If a specific vendor is named, hit their official site directly.
    for dom in _vendor_in_question(q):
        qs.append(f"{q} site:{dom}")
    out: list[Chunk] = []
    seen: set[str] = set()
    for qq in qs[:8]:
        try:
            chunks = web_search(qq, max_results=max_per_query)
        except Exception as e:
            log.warning("deep_web_research '%s' failed: %s", qq, e)
            continue
        for c in chunks:
            sig = (c.text or "")[:200]
            if sig and sig not in seen:
                seen.add(sig)
                out.append(c)
    return out


# --- agentic research planner ------------------------------------------
_PLANNER_SYS = (
    "You are an expert IoT-wireless research planner. Given a USER QUESTION and a "
    "PEEK at what local docs already say, decide what live research is needed to "
    "answer it accurately and completely. You MUST return STRICT JSON only.\n"
    "Authoritative domains you should prefer: bluetooth.com, ieee.org, csa-iot.org, "
    "threadgroup.org, wi-fi.org, fira.org, infineon.com, nxp.com, silabs.com, "
    "nordicsemi.com, qualcomm.com, ti.com, renesas.com, st.com, espressif.com, "
    "mediatek.com, realtek.com, synaptics.com.\n"
    "Use site: operators in queries when the question mentions a specific company, "
    "standard, spec version, product code, or technical detail (PHY, modulation, "
    "data rate, power class, throughput, etc.).\n"
    "Be aggressive: if the question is technical or asks to teach/explain/compare, "
    "ALWAYS plan at least 4-6 targeted queries and 0-3 specific URLs to fetch.\n"
    "CRITICAL: when the question mentions 'competitors' / 'competition' / 'rivals' / "
    "'industry' WITHOUT naming a specific vendor, you MUST enumerate AT LEAST these "
    "vendors and emit one query per vendor: Nordic Semiconductor, Silicon Labs, NXP, "
    "STMicroelectronics, Qualcomm, Synaptics, Renesas, Texas Instruments, Espressif, "
    "MediaTek, Realtek. Do NOT return a single generic 'competitors' query \u2014 fan out. "
    "NEVER include Infineon itself in competitor queries \u2014 the user works AT Infineon, "
    "they want news about OTHER vendors. Strip 'Infineon', 'AIROC', 'our' from the query.\n"
    "For news/recent/latest questions, include the current year (and last year) in queries "
    "and add 'news' / 'announcement' / 'press release' / 'product launch' keywords."
)

_PLANNER_SCHEMA = """
Return ONLY this JSON shape (no prose, no markdown):
{
  "needs_web": true,
  "queries": ["query 1 with site:domain.com if helpful", "query 2", "..."],
  "urls": ["https://exact.url/to/fetch", "..."],
  "rationale": "one short sentence on why"
}
- needs_web=false ONLY if the question is purely conversational ("hi", "thanks").
- queries: 3-12 items, concrete, include site: filters and key technical terms (PHY,
  Mbps, modulation, version numbers) when relevant. Include EVERY spec aspect the
  question implies (e.g. for 'Bluetooth HDT' include both BR/EDR rates AND LE rates
  AND Bluetooth 6.x release notes).
- For multi-vendor questions you MUST emit one query per vendor (see system prompt).
- urls: only include if you know an exact official spec / product page URL.
"""


def _competitor_list() -> list[str]:
    """Load competitor names from data/competitors.json + known fallback list."""
    base = ["Nordic Semiconductor", "Silicon Labs", "NXP", "STMicroelectronics",
            "Qualcomm", "Synaptics", "Renesas", "Texas Instruments",
            "Espressif", "MediaTek", "Realtek"]
    try:
        p = ROOT / "data" / "competitors.json"
        if p.exists():
            data = json.loads(p.read_text(encoding="utf-8"))
            extras = []
            if isinstance(data, list):
                for item in data:
                    if isinstance(item, str):
                        extras.append(item)
                    elif isinstance(item, dict):
                        n = item.get("name") or item.get("vendor") or item.get("company")
                        if n:
                            extras.append(str(n))
            elif isinstance(data, dict):
                for v in data.values():
                    if isinstance(v, str):
                        extras.append(v)
                    elif isinstance(v, dict):
                        n = v.get("name") or v.get("vendor")
                        if n:
                            extras.append(str(n))
            # merge, dedupe (case-insensitive)
            seen = {b.lower() for b in base}
            for e in extras:
                if e.lower() not in seen:
                    base.append(e); seen.add(e.lower())
    except Exception as e:
        log.warning("competitor list load failed: %s", e)
    return base


_GENERIC_COMPETITOR_TRIGGERS = ("competitor", "competitors", "competition",
                                "rivals", "industry players", "other vendors",
                                "all vendors")


def _expand_competitor_queries(question: str) -> list[str]:
    """If the question is about competitors generically, fan out one query per vendor.
    Excludes Infineon itself (we want competitor news, not our own)."""
    s = (question or "").lower()
    if not any(t in s for t in _GENERIC_COMPETITOR_TRIGGERS):
        return []
    # If a SPECIFIC competitor was named (other than Infineon), the planner handles it.
    named = [d for d in _vendor_in_question(question) if d != "infineon.com"]
    if named:
        return []
    year = time.strftime("%Y")
    last_year = str(int(year) - 1)
    out = []
    base_intent = re.sub(r"\b(infineon(?:'s)?|airoc|our|"
                         r"competitor(?:'s|s)?|competition|rivals|"
                         r"industry players|other vendors|all vendors)\b",
                         "", question, flags=re.I).strip(" ?.,;:")
    if not base_intent:
        base_intent = "wireless IoT product news announcement"
    for v in _competitor_list():
        if v.lower() == "infineon" or "infineon" in v.lower():
            continue  # skip Infineon itself
        out.append(f"{v} {base_intent} {year} {last_year}")
    return out


def _retriever_peek(question: str, retriever: Retriever | None) -> str:
    """Cheap snapshot of what local docs say so the planner can spot gaps."""
    try:
        retriever = retriever or Retriever()
        hits = retriever.search(question, k=4)
    except Exception:
        return "(no local index)"
    if not hits:
        return "(local index returned 0 hits)"
    parts = []
    for c, _s in hits:
        parts.append(f"- {c.title or c.source}: {(c.text or '').strip()[:240]}")
    return "\n".join(parts)


def _plan_research(question: str, retriever: Retriever | None) -> dict:
    """Ask an LLM what to search/fetch. Returns {needs_web, queries, urls, rationale}.
    Falls back to a heuristic plan if no LLM available."""
    api_key = os.environ.get("GEMINI_API_KEY") or os.environ.get("GOOGLE_API_KEY")
    if not api_key:
        return _heuristic_plan(question)
    try:
        import google.generativeai as genai
        genai.configure(api_key=api_key)
        peek = _retriever_peek(question, retriever)
        prompt = (
            f"USER QUESTION:\n{question}\n\n"
            f"WHAT LOCAL DOCS SAY (top-4 snippets):\n{peek}\n\n"
            f"TASK: plan the live research needed.\n{_PLANNER_SCHEMA}"
        )
        # Use the cheapest models for planning to save quota.
        for name in ("gemini-2.5-flash-lite", "gemini-2.0-flash-lite",
                     "gemini-flash-lite-latest"):
            try:
                m = genai.GenerativeModel(name, system_instruction=_PLANNER_SYS)
                resp = m.generate_content(
                    prompt,
                    generation_config={"response_mime_type": "application/json",
                                       "temperature": 0.1},
                )
                text = (getattr(resp, "text", None) or "").strip()
                if not text:
                    continue
                text = re.sub(r"^```(?:json)?\s*|\s*```$", "", text, flags=re.M)
                if not text.startswith("{"):
                    mm = re.search(r"\{.*\}", text, flags=re.S)
                    if mm:
                        text = mm.group(0)
                plan = json.loads(text)
                # sanitize
                plan["queries"] = [str(q).strip() for q in (plan.get("queries") or []) if str(q).strip()][:12]
                plan["urls"] = [str(u).strip() for u in (plan.get("urls") or []) if str(u).strip().startswith("http")][:5]
                plan["needs_web"] = bool(plan.get("needs_web", True))
                # Safety net: if user asked about competitors generically, force per-vendor queries.
                extra = _expand_competitor_queries(question)
                if extra:
                    seen = {q.lower() for q in plan["queries"]}
                    for q in extra:
                        if q.lower() not in seen:
                            plan["queries"].append(q); seen.add(q.lower())
                    plan["queries"] = plan["queries"][:14]
                    plan["needs_web"] = True
                return plan
            except Exception as e:
                msg = str(e).lower()
                if any(k in msg for k in ("429", "quota", "rate", "404", "500", "503")):
                    continue
                log.warning("planner %s failed: %s", name, e)
                continue
    except Exception as e:
        log.warning("planner setup failed: %s", e)
    return _heuristic_plan(question)


def _heuristic_plan(question: str) -> dict:
    """Used when planner LLM unavailable. Falls back to keyword heuristics."""
    if not question.strip():
        return {"needs_web": False, "queries": [], "urls": [], "rationale": "empty"}
    # Generic-competitor expansion ALWAYS runs (cheap & deterministic).
    expanded = _expand_competitor_queries(question)
    spec_q = _is_spec_question(question)
    if not spec_q and not expanded:
        return {"needs_web": True, "queries": [question], "urls": [],
                "rationale": "default web search"}
    qs = [question]
    if spec_q:
        qs += [f"{question} specification 2025 2026",
               f"{question} PHY data rate Mbps modulation"]
    ql = question.lower()
    if "bluetooth" in ql or "ble" in ql or "hdt" in ql or "auracast" in ql or "le audio" in ql:
        qs.append(f"{question} site:bluetooth.com")
    if "thread" in ql or "15.4" in ql:
        qs.append(f"{question} site:threadgroup.org")
    if "matter" in ql:
        qs.append(f"{question} site:csa-iot.org")
    if "wi-fi" in ql or "wifi" in ql or "802.11" in ql:
        qs.append(f"{question} site:wi-fi.org")
    for dom in _vendor_in_question(question):
        qs.append(f"{question} site:{dom}")
    qs.extend(expanded)
    # dedupe preserve order
    seen = set(); ded = []
    for q in qs:
        if q.lower() not in seen:
            seen.add(q.lower()); ded.append(q)
    return {"needs_web": True, "queries": ded[:14], "urls": [], "rationale": "heuristic"}


def _execute_plan(plan: dict, max_per_query: int = 3) -> list[Chunk]:
    """Run all queries + URL fetches from the plan, dedupe, return chunks."""
    out: list[Chunk] = []
    seen: set[str] = set()

    def _add(chunks):
        for c in chunks:
            sig = (c.text or "")[:200]
            if sig and sig not in seen:
                seen.add(sig)
                out.append(c)

    for u in plan.get("urls", []):
        try:
            _add(fetch_url(u))
        except Exception as e:
            log.warning("plan fetch_url %s: %s", u, e)
    for q in plan.get("queries", []):
        try:
            _add(web_search(q, max_results=max_per_query))
        except Exception as e:
            log.warning("plan web_search %r: %s", q, e)
    return out


def ask(question: str, history: list[dict] | None = None,
        retriever: Retriever | None = None,
        allow_web: bool = True) -> Answer:
    retriever = retriever or Retriever()
    hits = retriever.search(question, k=12)
    used_web = False
    plan: dict = {}

    # Always pre-fetch any inline URLs the user pasted.
    urls_in_q = re.findall(r"https?://[^\s)\]\"']+", question)
    if allow_web and urls_in_q:
        for u in urls_in_q[:3]:
            try:
                fetched = fetch_url(u)
                if fetched:
                    used_web = True
                    ingest_web_results(fetched)
            except Exception as e:
                log.warning("inline url fetch failed: %s", e)
        if used_web:
            retriever.load()
            hits = retriever.search(question, k=12) or hits

    # AGENTIC: let the LLM plan its own research, then execute it.
    if allow_web:
        plan = _plan_research(question, retriever)
        if plan.get("needs_web"):
            # When fanning out across many vendors, keep per-query results modest
            # so total stays bounded; otherwise allow more depth per query.
            per_q = 2 if len(plan.get("queries", [])) >= 8 else 3
            web_chunks = _execute_plan(plan, max_per_query=per_q)
            if web_chunks:
                used_web = True
                added = ingest_web_results(web_chunks)
                if added:
                    retriever.load()
                    hits = retriever.search(question, k=18) or hits
                else:
                    hits = (hits or []) + [(c, 0.5) for c in web_chunks[:12]]

    # Boost authoritative-domain hits.
    if hits:
        def _boost(item):
            c, s = item
            src = (c.source or "").lower()
            b = s
            if c.kind == "report":  # curated, human-facing site content
                b += 0.8
            if any(d in src for d in _AUTHORITATIVE_DOMAINS):
                b += 1.0
            return (c, b)
        hits = sorted([_boost(h) for h in hits], key=lambda x: x[1], reverse=True)[:12]

    # When the question is about competitors generically, drop Infineon-sourced hits
    # so the LLM can't fall back to talking about our own products.
    if any(t in question.lower() for t in _GENERIC_COMPETITOR_TRIGGERS) \
       and not [d for d in _vendor_in_question(question) if d != "infineon.com"]:
        def _is_infineon(c):
            src = (c.source or "").lower()
            title = (c.title or "").lower()
            return ("infineon" in src or "airoc" in src
                    or "infineon" in title or "airoc" in title)
        filtered = [(c, s) for c, s in hits if not _is_infineon(c)]
        if filtered:
            hits = filtered
            log.info("ask: filtered out %d Infineon-sourced hits for competitor question",
                     len(filtered))

    context, cites = _format_context(hits)

    models = _gemini_models()
    if not models:
        # Retrieval-only mode
        if not hits:
            return Answer(reply=("No indexed knowledge yet and no LLM key set. "
                                 "Run the AI Reindex button (or `python run.py ai-index`), "
                                 "and set GEMINI_API_KEY for full answers."),
                          sources=[], used_web=used_web, backend="retrieval-only")
        body = _synthesize_from_hits(question, hits, cites)
        return Answer(reply=("No GEMINI_API_KEY set — synthesized answer from retrieved snippets:\n\n" + body),
                      sources=cites, used_web=used_web, backend="retrieval-only")

    # Build prompt with history.
    hist_text = ""
    for turn in (history or [])[-6:]:
        role = turn.get("role", "user")
        content = turn.get("content", "")
        hist_text += f"\n[{role.upper()}] {content}"

    # Detect competitor-focused question (no specific competitor named).
    is_competitor_q = (
        any(t in question.lower() for t in _GENERIC_COMPETITOR_TRIGGERS)
        and not [d for d in _vendor_in_question(question) if d != "infineon.com"]
    )

    prompt = (
        f"CONTEXT (numbered, cite as [#n]):\n{context or '(no indexed context)'}\n\n"
        f"CONVERSATION SO FAR:{hist_text}\n\n"
        f"USER QUESTION: {question}\n\n"
        f"Answer following the rules in your system instruction.\n\n"
        f"FORBIDDEN PHRASES (the runtime fetched live web pages for you, listed in CONTEXT "
        f"above as items marked '(web)' — USE THEM): 'I do not have access', "
        f"'I cannot browse', 'I would need to perform searches', 'I would need to "
        f"consult', 'the provided context does not contain', 'to get specific X I would'. "
        f"If a vendor or topic isn't covered in CONTEXT, write a SHORT note like "
        f"'(no fresh data on Vendor X this round; ask me to dig deeper)' and continue "
        f"with what IS covered. Do not refuse."
    )
    if is_competitor_q:
        prompt += (
            "\n\nIMPORTANT \u2014 the user is asking about Infineon's COMPETITORS. "
            "Do NOT discuss Infineon, AIROC, or our own product roadmap. "
            "Focus ENTIRELY on the other vendors (Nordic, Silicon Labs, NXP, "
            "STMicroelectronics, Qualcomm, Synaptics, Renesas, Texas Instruments, "
            "Espressif, MediaTek, Realtek, etc.). Structure the answer with one "
            "section per competitor that has news, each with cited bullet points. "
            "Skip vendors with no fresh data \u2014 do not list them as 'no news'."
        )

    text, used_model, last_err = _run_model_chain(models, prompt)

    # AGENTIC SELF-CRITIQUE: if the answer admits gaps, do another research round
    # automatically and re-answer.
    if text and allow_web and _answer_has_gaps(text):
        gap_queries = _gap_queries(text, question)
        log.info("ask: self-critique detected gaps -> %d follow-up queries", len(gap_queries))
        more_chunks: list[Chunk] = []
        for q in gap_queries[:5]:
            try:
                more_chunks.extend(web_search(q, max_results=2))
            except Exception:
                pass
        if more_chunks:
            ingest_web_results(more_chunks)
            retriever.load()
            hits2 = retriever.search(question, k=14) or hits
            if hits2:
                hits2 = sorted(
                    [((c, s + 1.0) if any(d in (c.source or "").lower() for d in _AUTHORITATIVE_DOMAINS) else (c, s))
                     for c, s in hits2],
                    key=lambda x: x[1], reverse=True)[:12]
            context2, cites2 = _format_context(hits2)
            prompt2 = (
                f"CONTEXT (numbered, cite as [#n]):\n{context2 or '(no indexed context)'}\n\n"
                f"CONVERSATION SO FAR:{hist_text}\n\n"
                f"USER QUESTION: {question}\n\n"
                f"YOUR PREVIOUS ATTEMPT FLAGGED GAPS. Use the freshly fetched CONTEXT "
                f"above to give a complete answer. Never say you cannot browse the web; "
                f"the context items marked '(web)' ARE live web content fetched for you."
            )
            text2, used_model2, _ = _run_model_chain(models, prompt2)
            if text2:
                return Answer(reply=text2, sources=cites2, used_web=True,
                              backend=f"gemini:{used_model2} (self-refined)")

    if text:
        return Answer(reply=text, sources=cites, used_web=used_web,
                      backend=f"gemini:{used_model}")

    # All models failed → escalate: pull more web context and synthesize without LLM.
    if allow_web and not used_web:
        web_chunks = web_search(question, max_results=6)
        if web_chunks:
            used_web = True
            ingest_web_results(web_chunks)
            retriever.load()
            hits = retriever.search(question, k=8) or hits
            _, cites = _format_context(hits)

    body = _synthesize_from_hits(question, hits, cites)
    note = (f"_LLM unavailable (all Gemini fallbacks failed: {str(last_err)[:160]}). "
            f"Showing a synthesized answer from your indexed docs"
            f"{' + fresh web search' if used_web else ''}._\n\n")
    return Answer(reply=note + body, sources=cites, used_web=used_web, backend="retrieval-fallback")


# --- model chain runner & self-critique helpers ------------------------
def _run_model_chain(models, prompt: str) -> tuple[str, str, str]:
    """Try each model until one returns text. Returns (text, model_name, last_err)."""
    last_err = ""
    for m in models:
        name = getattr(m, "model_name", "?")
        try:
            resp = m.generate_content(prompt)
            text = (getattr(resp, "text", None) or "").strip()
            if text:
                return text, name, ""
            last_err = "empty response"
        except Exception as e:
            msg = str(e)
            last_err = msg
            if any(k in msg.lower() for k in ("429", "quota", "rate", "500", "503",
                                              "unavailable", "404", "not found")):
                log.warning("Gemini model %s failed (%s) — falling back", name, msg.splitlines()[0][:120])
                continue
            log.exception("Gemini call failed on %s", name)
            continue
    return "", "", last_err


_GAP_PHRASES = (
    "i do not have", "i don't have", "i cannot", "i can't",
    "as an ai", "i'm unable", "i am unable",
    "not able to browse", "no access to", "without browsing",
    "cannot access", "do not have access", "don't have access",
    "verify on", "(verify ", "verify with",
    "outside my training", "knowledge cutoff", "i don't know",
    "unable to find", "unclear from the provided",
    "the provided context does not", "context does not contain",
    "would need to consult", "would need to check",
    "would need to perform", "would need to search",
    "would need targeted", "i would need to",
    "the provided context does not offer",
    "does not contain specific",
    "to get specific", "what's missing",
    "more information would be needed",
    "no specific information", "no specific news",
    "without further information",
)


def _answer_has_gaps(text: str) -> bool:
    s = (text or "").lower()
    return any(p in s for p in _GAP_PHRASES)


def _gap_queries(answer_text: str, original_q: str) -> list[str]:
    """Heuristically build follow-up queries from gap markers in the answer."""
    queries: list[str] = []
    s = answer_text or ""
    # Pull sentences that contain "verify on X" or "would need to check X"
    for m in re.finditer(r"(?:verify on|would need to check|consult)\s+([A-Za-z0-9\.\-/ ]+)",
                          s, flags=re.I):
        target = m.group(1).strip(" .,;:")[:80]
        if target:
            queries.append(f"{original_q} {target}")
    # Pull sentences containing "I (don't|cannot) ... about X"
    for m in re.finditer(r"(?:don't have|cannot find|do not have).{0,40}?(?:about|regarding|on)\s+([^\.\n]{4,80})",
                          s, flags=re.I):
        target = m.group(1).strip(" .,;:")[:80]
        if target:
            queries.append(f"{original_q} {target}")
    # Generic fallback queries on authoritative domains
    ql = original_q.lower()
    if "bluetooth" in ql or "ble" in ql:
        queries.append(f"{original_q} site:bluetooth.com")
    if "thread" in ql or "15.4" in ql:
        queries.append(f"{original_q} site:threadgroup.org")
    if "matter" in ql:
        queries.append(f"{original_q} site:csa-iot.org")
    for dom in _vendor_in_question(original_q):
        queries.append(f"{original_q} site:{dom}")
    if not queries:
        queries.append(f"{original_q} datasheet specification")
    # dedupe preserving order
    seen = set(); out = []
    for q in queries:
        if q not in seen:
            seen.add(q); out.append(q)
    return out


def _synthesize_from_hits(question: str, hits: list, cites: list[dict]) -> str:
    """Build a readable, cited answer from retrieved chunks when the LLM is unavailable."""
    if not hits:
        return ("I could not find anything in the local index or via web search for this. "
                "Try rephrasing, or run `python run.py ai-enrich` to pull fresh material.")
    lines = [f"**Question:** {question}", "", "**Top findings (cited):**", ""]
    for i, (c, _s) in enumerate(hits[:6], 1):
        title = cites[i - 1]["title"] if i - 1 < len(cites) else (c.title or c.source)
        kind = c.kind
        snippet = (c.text or "").strip().replace("\n", " ")
        if len(snippet) > 700:
            snippet = snippet[:700].rstrip() + "…"
        loc = f" (p.{c.page})" if c.kind == "doc" and c.page else ""
        lines.append(f"**[#{i}] {title}{loc}** _({kind})_")
        lines.append(snippet)
        lines.append("")
    lines.append("**Sources:**")
    for c in cites[:6]:
        lines.append(f"- [#{c['id']}] {c['title']} — `{c['source']}`")
    return "\n".join(lines)
